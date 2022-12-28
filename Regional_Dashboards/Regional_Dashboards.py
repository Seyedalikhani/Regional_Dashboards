
import matplotlib.pyplot as plt
import ftplib
import patoolib
import os
import glob
from zipfile import ZipFile
import pandas as pd
import pyodbc
import datetime
import jdatetime
import numpy as np
from pylab import figure, clf, plot, bar, stem, xlabel, ylabel, xlim, ylim, title, grid, axes, show, legend
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import cv2



# ****************************************************************************************************
# ((((((((((((((((((((((((((((((((( Connection to PERFORMANCEDB01 ))))))))))))))))))))))))))))))))))))
# ****************************************************************************************************
conn = pyodbc.connect('Driver={SQL Server};'
                      'Server=PERFORMANCEDB01;'
                      'Database=Performance_NAK;'
                      'Trusted_Connection=yes;')
conn_performanceDB = conn.cursor()



# /\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\
# ((((((((((((((((((((( Functions ))))))))))))))))))))))))
# \/\/\/\/\/\/\/\/\/\/\/\\/\/\/\/\/\/\/\/\/\/\/\\/\/\/\/\/

# This function is used to covert centimeter to inch
def cm_to_inch(value):
    return value/2.54


# This function is used to reduce sample rate with downsample Rate=Rate
def downsample(Vector,Rate):
    downsample_vec=[]
    first_Index=0
    downsample_vec.append(Vector[0])
    for k in range(len(Vector)):
        if (k-first_Index==Rate):
            downsample_vec.append(Vector[k])
            first_Index=k
    return downsample_vec










# ---------------------------------------------------------------------
# ----------------------- Technical Dashboard -------------------------
# ---------------------------------------------------------------------

# -----------------------  Volte


conn_performanceDB.execute("select Date, Contractor, [PIndex], sum([L1800 Traffic]+[L2600 Traffic]+[L2300 Traffic]+[L2100 Traffic]) as 'Volte'from Province_KPI_Score_Band_CS_Daily group by Date, Contractor, [PIndex] order by Date")
Volte_Table=conn_performanceDB.fetchall()


Volte_Hourly_KH=[]
VolteTime_Hourly_KH=[]
Volte_Hourly_KM=[]
VolteTime_Hourly_KM=[]
Volte_Hourly_YZ=[]
VolteTime_Hourly_YZ=[]
Volte_Hourly_CH=[]
VolteTime_Hourly_CH=[]
Volte_Hourly_SM=[]
VolteTime_Hourly_SM=[]

Volte_Hourly_TH=[]
VolteTime_Hourly_TH=[]

Volte_Hourly_KJ=[]
VolteTime_Hourly_KJ=[]

Volte_Hourly_MA=[]
VolteTime_Hourly_MA=[]
Volte_Hourly_GN=[]
VolteTime_Hourly_GN=[]
Volte_Hourly_GL=[]
VolteTime_Hourly_GL=[]

Volte_Hourly_KZ=[]
VolteTime_Hourly_KZ=[]
Volte_Hourly_AS=[]
VolteTime_Hourly_AS=[]
Volte_Hourly_AG=[]
VolteTime_Hourly_AG=[]

for i in range(len(Volte_Table)):

            Row_Data=str(Volte_Table[i])
            Row_Data=Row_Data.split(", ")
    
            Year1=Row_Data[0]
            Date=Year1[19:23]+"/"+Row_Data[1]+"/"+Row_Data[2]
            Contractor=Row_Data[5]
            Contractor=Contractor[1:len(Contractor)-1]
            Province_Index=Row_Data[6]
            Province_Index=Province_Index[1:len(Province_Index)-1]
            Volte=Row_Data[7]
            if Volte=='None)':
                Volte=0
            else:
                Volte=round(float(Volte[0:len(Volte)-1])/1000,3)
            if (Contractor=="NAK-Nokia" and Province_Index=='KH'):
                Volte_Hourly_KH.append(Volte)
                VolteTime_Hourly_KH.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='KM'):
                Volte_Hourly_KM.append(Volte)
                VolteTime_Hourly_KM.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='YZ'):
                Volte_Hourly_YZ.append(Volte)
                VolteTime_Hourly_YZ.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='CH'):
                Volte_Hourly_CH.append(Volte)
                VolteTime_Hourly_CH.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='SM'):
                Volte_Hourly_SM.append(Volte)
                VolteTime_Hourly_SM.append(Date)
            if (Contractor=="NAK-Tehran" and Province_Index=='TH'):
                Volte_Hourly_TH.append(Volte)
                VolteTime_Hourly_TH.append(Date)
            if (Contractor=="NAK-Alborz" and Province_Index=='KJ'):
                Volte_Hourly_KJ.append(Volte)
                VolteTime_Hourly_KJ.append(Date)
            if (Contractor=="NAK-North" and Province_Index=='MA'):
                Volte_Hourly_MA.append(Volte)
                VolteTime_Hourly_MA.append(Date)
            if (Contractor=="NAK-North" and Province_Index=='GL'):
                Volte_Hourly_GL.append(Volte)
                VolteTime_Hourly_GL.append(Date)
            if (Contractor=="NAK-North" and Province_Index=='GN'):
                Volte_Hourly_GN.append(Volte)
                VolteTime_Hourly_GN.append(Date)
            if (Contractor=="NAK-Huawei" and Province_Index=='KZ'):
                Volte_Hourly_KZ.append(Volte)
                VolteTime_Hourly_KZ.append(Date)
            if (Contractor=="NAK-Huawei" and Province_Index=='AS'):
                Volte_Hourly_AS.append(Volte)
                VolteTime_Hourly_AS.append(Date)
            if (Contractor=="NAK-Huawei" and Province_Index=='AG'):
                Volte_Hourly_AG.append(Volte)
                VolteTime_Hourly_AG.append(Date)

for t in range(5):
        if (t==0 ):
            Contractor="NAK-Alborz"

            downsample_Rate=round(len(VolteTime_Hourly_KJ)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(VolteTime_Hourly_KJ,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(VolteTime_Hourly_KJ, Volte_Hourly_KJ, label='Alborz')
            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" Volte Traffic (KErlang)", fontdict = font1)
            plt.ylim((0,60))
            grid(True)
            plt.savefig("Volte_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Volte_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            Volte_Cropped = image[x:w, y:h]
            cv2.imwrite("Volte_"+Contractor+".png", Volte_Cropped)

        if (t==1 ):
            Contractor="NAK-North"

            downsample_Rate=round(len(VolteTime_Hourly_MA)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(VolteTime_Hourly_MA,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(VolteTime_Hourly_KH, Volte_Hourly_MA, label='Mazandaran')
            ax1.plot(VolteTime_Hourly_KM, Volte_Hourly_GN, label='Golestan')
            ax1.plot(VolteTime_Hourly_YZ, Volte_Hourly_GL, label='Gilan')


            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" Volte Traffic (KErlang)", fontdict = font1)
            plt.ylim((0,40))
            grid(True)
            plt.savefig("Volte_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Volte_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            Volte_Cropped = image[x:w, y:h]
            cv2.imwrite("Volte_"+Contractor+".png", Volte_Cropped)

        if (t==2 ):
            Contractor="NAK-Tehran"

            downsample_Rate=round(len(VolteTime_Hourly_TH)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(VolteTime_Hourly_TH,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(VolteTime_Hourly_KH, Volte_Hourly_TH, label='Tehran')
            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" Volte Traffic (KErlang)", fontdict = font1)
            plt.ylim((0,300))
            grid(True)
            plt.savefig("Volte_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Volte_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            Volte_Cropped = image[x:w, y:h]
            cv2.imwrite("Volte_"+Contractor+".png", Volte_Cropped)


        if (t==3 ):
            Contractor="NAK-Huawei"

            downsample_Rate=round(len(VolteTime_Hourly_KZ)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(VolteTime_Hourly_KZ,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(VolteTime_Hourly_KH, Volte_Hourly_KZ, label='Khouzestan')
            ax1.plot(VolteTime_Hourly_KM, Volte_Hourly_AS, label='Azarbaijan Sharghi')
            ax1.plot(VolteTime_Hourly_YZ, Volte_Hourly_AG, label='Azarbaijan Gharbi')


            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" Volte Traffic (KErlang)", fontdict = font1)
            plt.ylim((0,40))
            grid(True)
            plt.savefig("Volte_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Volte_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            Volte_Cropped = image[x:w, y:h]
            cv2.imwrite("Volte_"+Contractor+".png", Volte_Cropped)

        if (t==4 ):
            Contractor="NAK-Nokia"

            downsample_Rate=round(len(VolteTime_Hourly_KH)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(VolteTime_Hourly_KH,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(VolteTime_Hourly_KH, Volte_Hourly_KH, label='Khorasan Razavi')
            ax1.plot(VolteTime_Hourly_KM, Volte_Hourly_KM, label='Kerman')
            ax1.plot(VolteTime_Hourly_YZ, Volte_Hourly_YZ, label='Yazd')
            ax1.plot(VolteTime_Hourly_CH, Volte_Hourly_CH, label='Chaharmahal')
            ax1.plot(VolteTime_Hourly_SM, Volte_Hourly_SM, label='Semnan')

            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" Volte Traffic (KErlang)", fontdict = font1)
            plt.ylim((0,50))
            grid(True)
            plt.savefig("Volte_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Volte_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            Volte_Cropped = image[x:w, y:h]
            cv2.imwrite("Volte_"+Contractor+".png", Volte_Cropped)




# -----------------------  L2100

conn_performanceDB.execute("select Day, Contractor, [Province Index], sum([L2100 Payload]) as 'L2100' from Province_KPI_Score_Band_PS_Daily group by Day, Contractor, [Province Index] order by Day")
L2100_Table=conn_performanceDB.fetchall()



L2100_Hourly_KH=[]
L2100Time_Hourly_KH=[]
L2100_Hourly_KM=[]
L2100Time_Hourly_KM=[]
L2100_Hourly_YZ=[]
L2100Time_Hourly_YZ=[]
L2100_Hourly_CH=[]
L2100Time_Hourly_CH=[]
L2100_Hourly_SM=[]
L2100Time_Hourly_SM=[]

L2100_Hourly_TH=[]
L2100Time_Hourly_TH=[]

L2100_Hourly_KJ=[]
L2100Time_Hourly_KJ=[]

L2100_Hourly_MA=[]
L2100Time_Hourly_MA=[]
L2100_Hourly_GN=[]
L2100Time_Hourly_GN=[]
L2100_Hourly_GL=[]
L2100Time_Hourly_GL=[]

L2100_Hourly_KZ=[]
L2100Time_Hourly_KZ=[]
L2100_Hourly_AS=[]
L2100Time_Hourly_AS=[]
L2100_Hourly_AG=[]
L2100Time_Hourly_AG=[]

for i in range(len(L2100_Table)):

            Row_Data=str(L2100_Table[i])
            Row_Data=Row_Data.split(", ")
    
            Year1=Row_Data[0]
            Date=Year1[19:23]+"/"+Row_Data[1]+"/"+Row_Data[2]
            Contractor=Row_Data[5]
            Contractor=Contractor[1:len(Contractor)-1]
            Province_Index=Row_Data[6]
            Province_Index=Province_Index[1:len(Province_Index)-1]
            L2100=Row_Data[7]
            if L2100=='None)':
                continue
            L2100=round(float(L2100[0:len(L2100)-1])/1000,3)
            if (Contractor=="NAK-Nokia" and Province_Index=='KH'):
                L2100_Hourly_KH.append(L2100)
                L2100Time_Hourly_KH.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='KM'):
                L2100_Hourly_KM.append(L2100)
                L2100Time_Hourly_KM.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='YZ'):
                L2100_Hourly_YZ.append(L2100)
                L2100Time_Hourly_YZ.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='CH'):
                L2100_Hourly_CH.append(L2100)
                L2100Time_Hourly_CH.append(Date)
            if (Contractor=="NAK-Nokia" and Province_Index=='SM'):
                L2100_Hourly_SM.append(L2100)
                L2100Time_Hourly_SM.append(Date)
            if (Contractor=="NAK-Tehran" and Province_Index=='TH'):
                L2100_Hourly_TH.append(L2100)
                L2100Time_Hourly_TH.append(Date)
            if (Contractor=="NAK-Alborz" and Province_Index=='KJ'):
                L2100_Hourly_KJ.append(L2100)
                L2100Time_Hourly_KJ.append(Date)
            if (Contractor=="NAK-North" and Province_Index=='MA'):
                L2100_Hourly_MA.append(L2100)
                L2100Time_Hourly_MA.append(Date)
            if (Contractor=="NAK-North" and Province_Index=='GL'):
                L2100_Hourly_GL.append(L2100)
                L2100Time_Hourly_GL.append(Date)
            if (Contractor=="NAK-North" and Province_Index=='GN'):
                L2100_Hourly_GN.append(L2100)
                L2100Time_Hourly_GN.append(Date)
            if (Contractor=="NAK-Huawei" and Province_Index=='KZ'):
                L2100_Hourly_KZ.append(L2100)
                L2100Time_Hourly_KZ.append(Date)
            if (Contractor=="NAK-Huawei" and Province_Index=='AS'):
                L2100_Hourly_AS.append(L2100)
                L2100Time_Hourly_AS.append(Date)
            if (Contractor=="NAK-Huawei" and Province_Index=='AG'):
                L2100_Hourly_AG.append(L2100)
                L2100Time_Hourly_AG.append(Date)

for t in range(8):
        if (t==0 ):
            Contractor="NAK-Alborz"
            downsample_Rate=round(len(L2100Time_Hourly_KJ)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(L2100Time_Hourly_KJ,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(L2100Time_Hourly_KJ, L2100_Hourly_KJ, label='Alborz')
            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig("L2100_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            L2100_Cropped = image[x:w, y:h]
            cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)
        if (t==1 ):
            Contractor="NAK-North"

            downsample_Rate=round(len(L2100Time_Hourly_MA)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(L2100Time_Hourly_MA,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(L2100Time_Hourly_MA, L2100_Hourly_MA, label='Mazandaran')
            ax1.plot(L2100Time_Hourly_GN, L2100_Hourly_GN, label='Golestan')
            ax1.plot(L2100Time_Hourly_GL, L2100_Hourly_GL, label='Gilan')

            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig("L2100_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            L2100_Cropped = image[x:w, y:h]
            cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)
        if (t==2 ):
            Contractor="NAK-Tehran"

            downsample_Rate=round(len(L2100Time_Hourly_TH)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(L2100Time_Hourly_TH,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(L2100Time_Hourly_KH, L2100_Hourly_TH, label='Tehran')
            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig("L2100_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            L2100_Cropped = image[x:w, y:h]
            cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)


        if (t==3 ):
            Contractor="NAK-Huawei"

            downsample_Rate=round(len(L2100Time_Hourly_KZ)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(L2100Time_Hourly_KZ,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(L2100Time_Hourly_KZ, L2100_Hourly_KZ, label='Khouzestan')
            ax1.plot(L2100Time_Hourly_AS, L2100_Hourly_AS, label='Azarbaijan Sharghi')
            ax1.plot(L2100Time_Hourly_AG, L2100_Hourly_AG, label='Azarbaijan Gharbi')

            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig("L2100_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            L2100_Cropped = image[x:w, y:h]
            cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)

        if (t==4 ):
            Contractor="NAK-Nokia"

            downsample_Rate=round(len(L2100Time_Hourly_KH)/50)
            fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            x_Downsample=downsample(L2100Time_Hourly_KH,downsample_Rate)
            X_Vec=[]
            x_index=0
            while len(X_Vec)!=len(x_Downsample):
                 X_Vec.append(x_index)
                 x_index=x_index+downsample_Rate
            ax1.plot(L2100Time_Hourly_KH, L2100_Hourly_KH, label='Khorasan Razavi')
            ax1.plot(L2100Time_Hourly_KM, L2100_Hourly_KM, label='Kerman')
            ax1.plot(L2100Time_Hourly_YZ, L2100_Hourly_YZ, label='Yazd')
            ax1.plot(L2100Time_Hourly_CH, L2100_Hourly_CH, label='Chaharmahal')
            ax1.plot(L2100Time_Hourly_SM, L2100_Hourly_SM, label='Semnan')

            ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            leg = ax1.legend();
            font1 = {'family':'serif','color':'black','size':12}
            plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig("L2100_"+Contractor+".png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            y=100
            x=20
            h=1300
            w=1050
            L2100_Cropped = image[x:w, y:h]
            cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)

    


# -----------------------  CS PS Province



conn_performanceDB.execute("select Wk,PIndex, Contractor, sum([Total Voice Traffic (Erlang)]) as 'Total Voice Traffic (Erlang)' from ("+
                           "select Wk, Contractor,PIndex, avg([Total Voice Traffic (Erlang)]) as 'Total Voice Traffic (Erlang)' from "+
                            "Province_KPI_Score_Band_CS_Daily group by Wk, Contractor, PIndex ) tble group by Wk, Contractor, PIndex  order by Wk")
CS_Province_Table=conn_performanceDB.fetchall()


CS_Traffic_KH=[]
CS_Traffic_KM=[]
CS_Traffic_YZ=[]
CS_Traffic_CH=[]
CS_Traffic_SM=[]

CS_Traffic_KJ=[]
CS_Traffic_TH=[]
CS_Traffic_AS=[]
CS_Traffic_AG=[]
CS_Traffic_KZ=[]
CS_Traffic_MA=[]
CS_Traffic_GL=[]
CS_Traffic_GN=[]
# ----------------------- CS ----------------------------------

for i in range(len(CS_Province_Table)):


    Row_Data=str(CS_Province_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Province=Row_Data[1]
    Contractor=Row_Data[2]
    CS_Traffic=Row_Data[3]
    Week=Week[2:9]
    Province=Province[1:len(Province)-1]
    Contractor=Contractor[1:len(Contractor)-1]
    CS_Traffic_Val=round(float(CS_Traffic[0:len(CS_Traffic)-1])/1e3,3)

    #if Week=='1401-33':
    #    break



    if (Contractor=='NAK-Alborz'):
        CS_Traffic_KJ.append(CS_Traffic_Val)

    if (Contractor=='NAK-Tehran'):
        CS_Traffic_TH.append(CS_Traffic_Val)

    if (Contractor=='NAK-North'and Province=='MA'):
        CS_Traffic_MA.append(CS_Traffic_Val)
    if (Contractor=='NAK-North'and Province=='GN'):
        CS_Traffic_GN.append(CS_Traffic_Val)
    if (Contractor=='NAK-North'and Province=='GL'):
        CS_Traffic_GL.append(CS_Traffic_Val)

    if (Contractor=='NAK-Nokia' and Province=='KH'):
        CS_Traffic_KH.append(CS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='KM'):
        CS_Traffic_KM.append(CS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='YZ'):
        CS_Traffic_YZ.append(CS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='CH'):
        CS_Traffic_CH.append(CS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='SM'):
        CS_Traffic_SM.append(CS_Traffic_Val)
    if (Contractor=='NAK-Huawei'and Province=='KZ'):
        CS_Traffic_KZ.append(CS_Traffic_Val)
    if (Contractor=='NAK-Huawei'and Province=='AS'):
        CS_Traffic_AS.append(CS_Traffic_Val)
    if (Contractor=='NAK-Huawei'and Province=='AG'):
        CS_Traffic_AG.append(CS_Traffic_Val)





for t in range(8):
        #if (t==0 ):
        #    Contractor="NAK-Alborz"
        #    L2100_Hourly=L2100_Hourly_NAK_Alborz
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Alborz
        #if (t==1 ):
        #    Contractor="NAK-North"
        #    L2100_Hourly=L2100_Hourly_NAK_North
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_North
        if (t==2 ):
            Contractor="NAK-Tehran"

            #downsample_Rate=round(len(L2100Time_Hourly_TH)/50)
            #fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            #x_Downsample=downsample(L2100Time_Hourly_TH,downsample_Rate)
            #X_Vec=[]
            #x_index=0
            #while len(X_Vec)!=len(x_Downsample):
            #     X_Vec.append(x_index)
            #     x_index=x_index+downsample_Rate
            #ax1.plot(L2100Time_Hourly_KH, L2100_Hourly_TH, label='Tehran')
            #ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            #leg = ax1.legend();
            #font1 = {'family':'serif','color':'black','size':12}
            #plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            #grid(True)
            #plt.savefig("L2100_"+Contractor+".png")

            #image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            #y=100
            #x=20
            #h=1300
            #w=1050
            #L2100_Cropped = image[x:w, y:h]
            #cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)


        #if (t==3 ):
        #    Contractor="NAK-Huawei"
        #    L2100_Hourly=L2100_Hourly_NAK_Huawei
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Huawei
        if (t==4 ):
            Contractor="NAK-Nokia"


             # Sort Data Based on Last Values
            Last_CS_Traffic_NAK_Nokia_Value=[CS_Traffic_KH[len(CS_Traffic_KH)-1], CS_Traffic_KM[len(CS_Traffic_KM)-1], CS_Traffic_YZ[len(CS_Traffic_YZ)-1],  CS_Traffic_CH[len(CS_Traffic_CH)-1], CS_Traffic_SM[len(CS_Traffic_SM)-1]]
            Index_of_Sort=np.argsort(Last_CS_Traffic_NAK_Nokia_Value)

            Data_Sorted_Array=[]
            x_Labels=[];
            for k in range(len(Index_of_Sort)):
                if Index_of_Sort[k]==0:
                    Data_Sorted_Array.append(CS_Traffic_KH)
                    x_Labels.append('Khorasan Razavi')
                if Index_of_Sort[k]==1:
                    Data_Sorted_Array.append(CS_Traffic_KM)
                    x_Labels.append('Kerman')
                if Index_of_Sort[k]==2:
                    Data_Sorted_Array.append(CS_Traffic_YZ)
                    x_Labels.append('Yazd')
                if Index_of_Sort[k]==3:
                    Data_Sorted_Array.append(CS_Traffic_CH)
                    x_Labels.append('Chaharmahal')
                if Index_of_Sort[k]==4:
                    Data_Sorted_Array.append(CS_Traffic_SM)
                    x_Labels.append('Semnan')


            data=np.array(Data_Sorted_Array)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(12)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_CS_Traffic_NAK_Nokia_Value):
                plt.text( i + dx[31],Last_CS_Traffic_NAK_Nokia_Value[Index_of_Sort[i]] , str(Last_CS_Traffic_NAK_Nokia_Value[Index_of_Sort[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels)
            font1 = {'family':'serif','color':'black','size':15}
            plt.title(Contractor+" Total Traffic (KErlang)", fontdict = font1)
            grid(True)
            plt.savefig('CS_Traffic_Bar_'+Contractor+'.png')

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CS_Traffic_Bar_"+Contractor+".png")
            y=80
            x=20
            h=1000
            w=520
            CS_Traffic_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite('CS_Traffic_Bar_'+Contractor+'.png', CS_Traffic_Bar_Cropped)




            Data_Sorted_Pie=[Data_Sorted_Array[0][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[1][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[2][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[3][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[4][len(Data_Sorted_Array[0])-1]]
            plt.figure(figsize=(cm_to_inch(12),cm_to_inch(10)))
            #font1 = {'family':'serif','color':'black','size':15}
            #plt.title("Total Traffic (%)", fontdict = font1)
            plt.pie(Data_Sorted_Pie,labels =x_Labels,autopct='%1.1f%%')
            plt.savefig('CS_Traffic_Bar_'+Contractor+'_Percentage_1.png')


            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CS_Traffic_Bar_"+Contractor+"_Percentage_1.png")
            y=50
            x=20
            h=450
            w=350
            CS_Traffic_Bar_Percentage_1_Cropped = image[x:w, y:h]
            cv2.imwrite('CS_Traffic_Bar_'+Contractor+'_Percentage_1.png', CS_Traffic_Bar_Percentage_1_Cropped)



# -----------------------  PS PS Province



conn_performanceDB.execute("select Wk,[Province Index], Contractor, sum([Total Payload (GB)]) as 'Total Payload (GB)' from ("+
                           "select Wk, Contractor,[Province Index], avg([Total Payload (GB)]) as 'Total Payload (GB)' from "+
                            "Province_KPI_Score_Band_PS_Daily group by Wk, Contractor, [Province Index] ) tble group by Wk, Contractor, [Province Index]  order by Wk")
PS_Province_Table=conn_performanceDB.fetchall()


PS_Traffic_KH=[]
PS_Traffic_KM=[]
PS_Traffic_YZ=[]
PS_Traffic_CH=[]
PS_Traffic_SM=[]



# ----------------------- PS ----------------------------------

for i in range(len(PS_Province_Table)):


    Row_Data=str(PS_Province_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Province=Row_Data[1]
    Contractor=Row_Data[2]
    PS_Traffic=Row_Data[3]
    Week=Week[2:9]
    Province=Province[1:len(Province)-1]
    Contractor=Contractor[1:len(Contractor)-1]
    PS_Traffic_Val=round(float(PS_Traffic[0:len(PS_Traffic)-1])/1e3,3)

    #if Week=='1401-33':
    #    break



    #if (Contractor=='NAK-Alborz'):
    #    PS_Traffic_NAK_Alborz.append(PS_Traffic_Val)
    #if (Contractor=='NAK-Tehran'):
    #    PS_Traffic_NAK_Tehran.append(PS_Traffic_Val)
    #if (Contractor=='NAK-North'):
    #    PS_Traffic_NAK_North.append(PS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='KH'):
        PS_Traffic_KH.append(PS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='KM'):
        PS_Traffic_KM.append(PS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='YZ'):
        PS_Traffic_YZ.append(PS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='CH'):
        PS_Traffic_CH.append(PS_Traffic_Val)
    if (Contractor=='NAK-Nokia' and Province=='SM'):
        PS_Traffic_SM.append(PS_Traffic_Val)
    #if (Contractor=='NAK-Huawei'):
    #    PS_Traffic_NAK_Huawei.append(PS_Traffic_Val)
    #if (Contractor=='Farafan'):
    #    PS_Traffic_Farafan.append(PS_Traffic_Val)
    #if (Contractor=='BR-TEL'):
    #    PS_Traffic_BR_TEL.append(PS_Traffic_Val)
    #if (Contractor=='Huawei'):
    #    PS_Traffic_Huawei.append(PS_Traffic_Val)






for t in range(8):
        #if (t==0 ):
        #    Contractor="NAK-Alborz"
        #    L2100_Hourly=L2100_Hourly_NAK_Alborz
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Alborz
        #if (t==1 ):
        #    Contractor="NAK-North"
        #    L2100_Hourly=L2100_Hourly_NAK_North
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_North
        if (t==2 ):
            Contractor="NAK-Tehran"

            #downsample_Rate=round(len(L2100Time_Hourly_TH)/50)
            #fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            #x_Downsample=downsample(L2100Time_Hourly_TH,downsample_Rate)
            #X_Vec=[]
            #x_index=0
            #while len(X_Vec)!=len(x_Downsample):
            #     X_Vec.append(x_index)
            #     x_index=x_index+downsample_Rate
            #ax1.plot(L2100Time_Hourly_KH, L2100_Hourly_TH, label='Tehran')
            #ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            #leg = ax1.legend();
            #font1 = {'family':'serif','color':'black','size':12}
            #plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            #grid(True)
            #plt.savefig("L2100_"+Contractor+".png")

            #image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            #y=100
            #x=20
            #h=1300
            #w=1050
            #L2100_Cropped = image[x:w, y:h]
            #cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)


        #if (t==3 ):
        #    Contractor="NAK-Huawei"
        #    L2100_Hourly=L2100_Hourly_NAK_Huawei
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Huawei
        if (t==4 ):
            Contractor="NAK-Nokia"


             # Sort Data Based on Last Values
            Last_PS_Traffic_NAK_Nokia_Value=[PS_Traffic_KH[len(PS_Traffic_KH)-1], PS_Traffic_KM[len(PS_Traffic_KM)-1], PS_Traffic_YZ[len(PS_Traffic_YZ)-1],  PS_Traffic_CH[len(PS_Traffic_CH)-1], PS_Traffic_SM[len(PS_Traffic_SM)-1]]
            Index_of_Sort=np.argsort(Last_PS_Traffic_NAK_Nokia_Value)

            Data_Sorted_Array=[]
            x_Labels=[];
            for k in range(len(Index_of_Sort)):
                if Index_of_Sort[k]==0:
                    Data_Sorted_Array.append(PS_Traffic_KH)
                    x_Labels.append('Khorasan Razavi')
                if Index_of_Sort[k]==1:
                    Data_Sorted_Array.append(PS_Traffic_KM)
                    x_Labels.append('Kerman')
                if Index_of_Sort[k]==2:
                    Data_Sorted_Array.append(PS_Traffic_YZ)
                    x_Labels.append('Yazd')
                if Index_of_Sort[k]==3:
                    Data_Sorted_Array.append(PS_Traffic_CH)
                    x_Labels.append('Chaharmahal')
                if Index_of_Sort[k]==4:
                    Data_Sorted_Array.append(PS_Traffic_SM)
                    x_Labels.append('Semnan')


            data=np.array(Data_Sorted_Array)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(12)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_PS_Traffic_NAK_Nokia_Value):
                plt.text( i + dx[31],Last_PS_Traffic_NAK_Nokia_Value[Index_of_Sort[i]] , str(Last_PS_Traffic_NAK_Nokia_Value[Index_of_Sort[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels)
            font1 = {'family':'serif','color':'black','size':15}
            plt.title(Contractor+" Total Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig('PS_Traffic_Bar_'+Contractor+'.png')

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\PS_Traffic_Bar_"+Contractor+".png")
            y=80
            x=20
            h=1000
            w=520
            PS_Traffic_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite('PS_Traffic_Bar_'+Contractor+'.png', PS_Traffic_Bar_Cropped)




            Data_Sorted_Pie=[Data_Sorted_Array[0][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[1][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[2][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[3][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[4][len(Data_Sorted_Array[0])-1]]
            plt.figure(figsize=(cm_to_inch(12),cm_to_inch(10)))
            #font1 = {'family':'serif','color':'black','size':15}
            #plt.title("Total Traffic (%)", fontdict = font1)
            plt.pie(Data_Sorted_Pie,labels =x_Labels,autopct='%1.1f%%')
            plt.savefig('PS_Traffic_Bar_'+Contractor+'_Percentage_1.png')


            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\PS_Traffic_Bar_"+Contractor+"_Percentage_1.png")
            y=50
            x=20
            h=450
            w=350
            PS_Traffic_Bar_Percentage_1_Cropped = image[x:w, y:h]
            cv2.imwrite('PS_Traffic_Bar_'+Contractor+'_Percentage_1.png', PS_Traffic_Bar_Percentage_1_Cropped)






# -----------------------  CC and RD



conn_performanceDB.execute("select Wk,[PIndex], Contractor, avg([CC2]) as 'CC2', avg([CC3]) as 'CC3', avg([CC]) as 'CC', sum([2G_TCH_Traffic]) as '2G_TCH_Traffic', sum([3G_CS_Traffic]) as '3G_CS_Traffic', sum([4G_Volte_Traffic]) as '4G_Volte_Traffic',  sum([Total_Traffic]) as 'Total_Traffic' from ("+
                           "select Wk, Contractor,PIndex, SUM([2G TCH Traffic]*[CC2 (%)])/sum([2G TCH Traffic]) as 'CC2',"+
                              "SUM([3G_CS_Traffic]*[CC3 (%)])/sum([3G_CS_Traffic]) as 'CC3',"+
							  "SUM(([2G TCH Traffic]+[3G_CS_Traffic])*[CC (%)])/sum([2G TCH Traffic]+[3G_CS_Traffic]) as 'CC',"+
							  "avg([2G TCH Traffic]) as '2G_TCH_Traffic',"+
							  "avg([3G_CS_Traffic]) as '3G_CS_Traffic',"+
							  "avg([L1800 Traffic]+[L2600 Traffic]+[L2300 Traffic]+[L2100 Traffic]) as '4G_Volte_Traffic',"+
							  "avg([Total Voice Traffic (Erlang)]) as 'Total_Traffic' "+
							  "from  Province_KPI_Score_Band_CS_Daily  group by Wk, PIndex, Contractor, PIndex ) tble  group by Wk, PIndex, Contractor  order by Wk")
CC_Province_Table=conn_performanceDB.fetchall()


conn_performanceDB.execute("select Wk, [Province Index], Contractor, avg([RD2]) as 'RD2', avg([RD3]) as 'RD3', avg([RD4]) as 'RD4', avg([RD]) as 'RD', sum([2G_PS_Payload]) as '2G_PS_Payload', sum([3G_PS_Payload]) as '3G_PS_Payload', sum([4G_PS_Payload]) as '4G_PS_Payload', sum([Total_Payload]) as 'Total_Payload' from ("+
                             "select Wk, Contractor,[Province Index], SUM([2G PS Traffic (GB)]*[RD2 (%)])/sum([2G PS Traffic (GB)]) as 'RD2',"+
                              "SUM([3G Payload (GB)]*[RD3 (%)])/sum([3G Payload (GB)]) as 'RD3',"+
							  "SUM([4G Payload (GB)]*[RD4 (%)])/sum([4G Payload (GB)]) as 'RD4',"+
							  "SUM([Total Payload (GB)]*[RD (%)])/sum([Total Payload (GB)]) as 'RD',"+
							  "avg([2G PS Traffic (GB)]) as '2G_PS_Payload',"+
							  "avg([3G Payload (GB)]) as '3G_PS_Payload',"+
							  "avg([4G Payload (GB)]) as '4G_PS_Payload',"+
							  "avg([Total Payload (GB)]) as 'Total_Payload' "+
							  "from  Province_KPI_Score_Band_PS_Daily group by Wk, [Province Index], Contractor, [Province Index] ) tble group by Wk, [Province Index], Contractor  order by Wk")
RD_Province_Table=conn_performanceDB.fetchall()


CC2_NAK_KH=[]
CC2_NAK_KM=[]
CC2_NAK_YZ=[]
CC2_NAK_CH=[]
CC2_NAK_SM=[]

CC3_NAK_KH=[]
CC3_NAK_KM=[]
CC3_NAK_YZ=[]
CC3_NAK_CH=[]
CC3_NAK_SM=[]


CC_NAK_KH=[]
CC_NAK_KM=[]
CC_NAK_YZ=[]
CC_NAK_CH=[]
CC_NAK_SM=[]





CS_Traffic_2G_KH=[]
CS_Traffic_2G_KM=[]
CS_Traffic_2G_YZ=[]
CS_Traffic_2G_CH=[]
CS_Traffic_2G_SM=[]


CS_Traffic_3G_KH=[]
CS_Traffic_3G_KM=[]
CS_Traffic_3G_YZ=[]
CS_Traffic_3G_CH=[]
CS_Traffic_3G_SM=[]



CS_Traffic_4G_KH=[]
CS_Traffic_4G_KM=[]
CS_Traffic_4G_YZ=[]
CS_Traffic_4G_CH=[]
CS_Traffic_4G_SM=[]




Week_Vec=[]

for i in range(len(CC_Province_Table)):
    Row_Data=str(CC_Province_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Province=Row_Data[1]
    Contractor=Row_Data[2]
    CC2_Str=Row_Data[3]
    CC2_Val=round(float(CC2_Str[0:len(CC2_Str)-1]),2)
    CC3_Str=Row_Data[4]
    CC3_Val=round(float(CC3_Str[0:len(CC3_Str)-1]),2)
    CC_Str=Row_Data[5]
    CC_Val=round(float(CC_Str[0:len(CC_Str)-1]),2)
    CS_2G_Str=Row_Data[6]
    CS_2G_Val=round(float(CS_2G_Str[0:len(CS_2G_Str)-1]),2)
    CS_3G_Str=Row_Data[7]
    CS_3G_Val=round(float(CS_3G_Str[0:len(CS_3G_Str)-1]),2)
    CS_4G_Str=Row_Data[8]
    CS_4G_Val=round(float(CS_4G_Str[0:len(CS_4G_Str)-1]),2)
    CS_Str=Row_Data[9]
    CS_Val=round(float(CS_Str[0:len(CS_Str)-1]),2)

    Week=Week[2:9]
    Contractor=Contractor[1:len(Contractor)-1]
    Province=Province[1:len(Province)-1]

    #if Week=='1401-33':
    #    break

    #if (Contractor=='NAK-Alborz'):
    #    Week_Vec.append('W'+Week[5:7])
    #    CC2_NAK_Alborz.append(CC2_Val)
    #    CC3_NAK_Alborz.append(CC3_Val)
    #    CC_NAK_Alborz.append(CC_Val)
    #    CS_Traffic_2G_NAK_Alborz.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_NAK_Alborz.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_NAK_Alborz.append(CS_4G_Val/1e3)
    #if (Contractor=='NAK-Tehran'):
    #    CC2_NAK_Tehran.append(CC2_Val)
    #    CC3_NAK_Tehran.append(CC3_Val)
    #    CC_NAK_Tehran.append(CC_Val)
    #    CS_Traffic_2G_NAK_Tehran.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_NAK_Tehran.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_NAK_Tehran.append(CS_4G_Val/1e3)
    #if (Contractor=='NAK-North'):
    #    CC2_NAK_North.append(CC2_Val)
    #    CC3_NAK_North.append(CC3_Val)
    #    CC_NAK_North.append(CC_Val)
    #    CS_Traffic_2G_NAK_North.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_NAK_North.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_NAK_North.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-Nokia' and Province=='KH'):
        Week_Vec.append('W'+Week[5:7])
        CC2_NAK_KH.append(CC2_Val)
        CC3_NAK_KH.append(CC3_Val)
        CC_NAK_KH.append(CC_Val)
        CS_Traffic_2G_KH.append(CS_2G_Val)
        CS_Traffic_3G_KH.append(CS_3G_Val)
        CS_Traffic_4G_KH.append(CS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='KM'):
        CC2_NAK_KM.append(CC2_Val)
        CC3_NAK_KM.append(CC3_Val)
        CC_NAK_KM.append(CC_Val)
        CS_Traffic_2G_KM.append(CS_2G_Val)
        CS_Traffic_3G_KM.append(CS_3G_Val)
        CS_Traffic_4G_KM.append(CS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='YZ'):
        CC2_NAK_YZ.append(CC2_Val)
        CC3_NAK_YZ.append(CC3_Val)
        CC_NAK_YZ.append(CC_Val)
        CS_Traffic_2G_YZ.append(CS_2G_Val)
        CS_Traffic_3G_YZ.append(CS_3G_Val)
        CS_Traffic_4G_YZ.append(CS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='CH'):
        CC2_NAK_CH.append(CC2_Val)
        CC3_NAK_CH.append(CC3_Val)
        CC_NAK_CH.append(CC_Val)
        CS_Traffic_2G_CH.append(CS_2G_Val)
        CS_Traffic_3G_CH.append(CS_3G_Val)
        CS_Traffic_4G_CH.append(CS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='SM'):
        CC2_NAK_SM.append(CC2_Val)
        CC3_NAK_SM.append(CC3_Val)
        CC_NAK_SM.append(CC_Val)
        CS_Traffic_2G_SM.append(CS_2G_Val)
        CS_Traffic_3G_SM.append(CS_3G_Val)
        CS_Traffic_4G_SM.append(CS_4G_Val)
    #if (Contractor=='NAK-Huawei'):
    #    CC2_NAK_Huawei.append(CC2_Val)
    #    CC3_NAK_Huawei.append(CC3_Val)
    #    CC_NAK_Huawei.append(CC_Val)
    #    CS_Traffic_2G_NAK_Huawei.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_NAK_Huawei.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_NAK_Huawei.append(CS_4G_Val/1e3)
    #if (Contractor=='Farafan'):
    #    CC2_Farafan.append(CC2_Val)
    #    CC3_Farafan.append(CC3_Val)
    #    CC_Farafan.append(CC_Val)
    #    CS_Traffic_2G_Farafan.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_Farafan.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_Farafan.append(CS_4G_Val/1e3)
    #if (Contractor=='BR-TEL'):
    #    CC2_BR_TEL.append(CC2_Val)
    #    CC3_BR_TEL.append(CC3_Val)
    #    CC_BR_TEL.append(CC_Val)
    #    CS_Traffic_2G_BR_TEL.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_BR_TEL.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_BR_TEL.append(CS_4G_Val/1e3)
    #if (Contractor=='Huawei'):
    #    CC2_Huawei.append(CC2_Val)
    #    CC3_Huawei.append(CC3_Val)
    #    CC_Huawei.append(CC_Val)
    #    CS_Traffic_2G_Huawei.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_Huawei.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_Huawei.append(CS_4G_Val/1e3)
    #if (Contractor=='IRAN'):
    #    CC2_Iran.append(CC2_Val)
    #    CC3_Iran.append(CC3_Val)
    #    CC_Iran.append(CC_Val)
    #    CS_Traffic_2G_Iran.append(CS_2G_Val/1e3)
    #    CS_Traffic_3G_Iran.append(CS_3G_Val/1e3)
    #    CS_Traffic_4G_Iran.append(CS_4G_Val/1e3)






for t in range(8):
        #if (t==0 ):
        #    Contractor="NAK-Alborz"
        #    L2100_Hourly=L2100_Hourly_NAK_Alborz
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Alborz
        #if (t==1 ):
        #    Contractor="NAK-North"
        #    L2100_Hourly=L2100_Hourly_NAK_North
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_North
        if (t==2 ):
            Contractor="NAK-Tehran"

            #downsample_Rate=round(len(L2100Time_Hourly_TH)/50)
            #fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            #x_Downsample=downsample(L2100Time_Hourly_TH,downsample_Rate)
            #X_Vec=[]
            #x_index=0
            #while len(X_Vec)!=len(x_Downsample):
            #     X_Vec.append(x_index)
            #     x_index=x_index+downsample_Rate
            #ax1.plot(L2100Time_Hourly_KH, L2100_Hourly_TH, label='Tehran')
            #ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            #leg = ax1.legend();
            #font1 = {'family':'serif','color':'black','size':12}
            #plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            #grid(True)
            #plt.savefig("L2100_"+Contractor+".png")

            #image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            #y=100
            #x=20
            #h=1300
            #w=1050
            #L2100_Cropped = image[x:w, y:h]
            #cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)


        #if (t==3 ):
        #    Contractor="NAK-Huawei"
        #    L2100_Hourly=L2100_Hourly_NAK_Huawei
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Huawei
        if (t==4 ):
            Contractor="NAK-Nokia"

            Last_CC2_NAK_Nokia=[CC2_NAK_KH[len(CC2_NAK_KH)-1], CC2_NAK_KM[len(CC2_NAK_KM)-1], CC2_NAK_YZ[len(CC2_NAK_YZ)-1], CC2_NAK_CH[len(CC2_NAK_CH)-1], CC2_NAK_SM[len(CC2_NAK_SM)-1]]
            Index_of_Sort_CC2=np.argsort(Last_CC2_NAK_Nokia)
            Data_Sorted_Array_CC2=[]
            x_Labels_CC2=[];
            for k in range(len(Index_of_Sort_CC2)):
                if Index_of_Sort_CC2[k]==0:
                    Data_Sorted_Array_CC2.append(CC2_NAK_KH)
                    x_Labels_CC2.append('Khorasan Razavi')
                if Index_of_Sort_CC2[k]==1:
                    Data_Sorted_Array_CC2.append(CC2_NAK_KM)
                    x_Labels_CC2.append('Kerman')
                if Index_of_Sort_CC2[k]==2:
                    Data_Sorted_Array_CC2.append(CC2_NAK_YZ)
                    x_Labels_CC2.append('Yazd')
                if Index_of_Sort_CC2[k]==3:
                    Data_Sorted_Array_CC2.append(CC2_NAK_CH)
                    x_Labels_CC2.append('Chaharmahal')
                if Index_of_Sort_CC2[k]==4:
                    Data_Sorted_Array_CC2.append(CC2_NAK_SM)
                    x_Labels_CC2.append('Semnan')


            data=np.array(Data_Sorted_Array_CC2)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_CC2_NAK_Nokia):
                plt.text( i + dx[31],Last_CC2_NAK_Nokia[Index_of_Sort_CC2[i]] , str(Last_CC2_NAK_Nokia[Index_of_Sort_CC2[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels_CC2)
            font1 = {'family':'serif','color':'black','size':17}
            plt.title("CC2(%)", fontdict = font1)
            plt.ylim(75, 100)
            grid(True)
            plt.savefig("CC2_"+Contractor+"_Provinces.png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC2_"+Contractor+"_Provinces.png")
            y=80
            x=10
            h=1000
            w=520
            CC2_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite("CC2_"+Contractor+"_Provinces.png", CC2_Bar_Cropped)





            Last_CC3_NAK_Nokia=[CC3_NAK_KH[len(CC3_NAK_KH)-1], CC3_NAK_KM[len(CC3_NAK_KM)-1], CC3_NAK_YZ[len(CC3_NAK_YZ)-1], CC3_NAK_CH[len(CC3_NAK_CH)-1], CC3_NAK_SM[len(CC3_NAK_SM)-1]]
            Index_of_Sort_CC3=np.argsort(Last_CC3_NAK_Nokia)
            Data_Sorted_Array_CC3=[]
            x_Labels_CC3=[];
            for k in range(len(Index_of_Sort_CC3)):
                if Index_of_Sort_CC3[k]==0:
                    Data_Sorted_Array_CC3.append(CC3_NAK_KH)
                    x_Labels_CC3.append('Khorasan Razavi')
                if Index_of_Sort_CC3[k]==1:
                    Data_Sorted_Array_CC3.append(CC3_NAK_KM)
                    x_Labels_CC3.append('Kerman')
                if Index_of_Sort_CC3[k]==2:
                    Data_Sorted_Array_CC3.append(CC3_NAK_YZ)
                    x_Labels_CC3.append('Yazd')
                if Index_of_Sort_CC3[k]==3:
                    Data_Sorted_Array_CC3.append(CC3_NAK_CH)
                    x_Labels_CC3.append('Chaharmahal')
                if Index_of_Sort_CC3[k]==4:
                    Data_Sorted_Array_CC3.append(CC3_NAK_SM)
                    x_Labels_CC3.append('Semnan')


            data=np.array(Data_Sorted_Array_CC3)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_CC3_NAK_Nokia):
                plt.text( i + dx[31],Last_CC3_NAK_Nokia[Index_of_Sort_CC3[i]] , str(Last_CC3_NAK_Nokia[Index_of_Sort_CC3[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels_CC3)
            font1 = {'family':'serif','color':'black','size':17}
            plt.title("CC3(%)", fontdict = font1)
            plt.ylim(75, 100)
            grid(True)
            plt.savefig("CC3_"+Contractor+"_Provinces.png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC3_"+Contractor+"_Provinces.png")
            y=80
            x=10
            h=1000
            w=520
            CC3_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite("CC3_"+Contractor+"_Provinces.png", CC3_Bar_Cropped)



            Last_CC_NAK_Nokia=[CC_NAK_KH[len(CC_NAK_KH)-1], CC_NAK_KM[len(CC_NAK_KM)-1], CC_NAK_YZ[len(CC_NAK_YZ)-1], CC_NAK_CH[len(CC_NAK_CH)-1], CC_NAK_SM[len(CC_NAK_SM)-1]]
            Index_of_Sort_CC=np.argsort(Last_CC_NAK_Nokia)
            Data_Sorted_Array_CC=[]
            x_Labels_CC=[];
            for k in range(len(Index_of_Sort_CC)):
                if Index_of_Sort_CC[k]==0:
                    Data_Sorted_Array_CC.append(CC_NAK_KH)
                    x_Labels_CC.append('Khorasan Razavi')
                if Index_of_Sort_CC[k]==1:
                    Data_Sorted_Array_CC.append(CC_NAK_KM)
                    x_Labels_CC.append('Kerman')
                if Index_of_Sort_CC[k]==2:
                    Data_Sorted_Array_CC.append(CC_NAK_YZ)
                    x_Labels_CC.append('Yazd')
                if Index_of_Sort_CC[k]==3:
                    Data_Sorted_Array_CC.append(CC_NAK_CH)
                    x_Labels_CC.append('Chaharmahal')
                if Index_of_Sort_CC[k]==4:
                    Data_Sorted_Array_CC.append(CC_NAK_SM)
                    x_Labels_CC.append('Semnan')


            data=np.array(Data_Sorted_Array_CC)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_CC_NAK_Nokia):
                plt.text( i + dx[31],Last_CC_NAK_Nokia[Index_of_Sort_CC[i]] , str(Last_CC_NAK_Nokia[Index_of_Sort_CC[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels_CC)
            font1 = {'family':'serif','color':'black','size':17}
            plt.title("CC(%)", fontdict = font1)
            plt.ylim(75, 100)
            grid(True)
            plt.savefig("CC_"+Contractor+"_Provinces.png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC_"+Contractor+"_Provinces.png")
            y=80
            x=10
            h=1000
            w=520
            CC_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite("CC_"+Contractor+"_Provinces.png", CC_Bar_Cropped)
 

RD2_NAK_KH=[]
RD2_NAK_KM=[]
RD2_NAK_YZ=[]
RD2_NAK_CH=[]
RD2_NAK_SM=[]

RD3_NAK_KH=[]
RD3_NAK_KM=[]
RD3_NAK_YZ=[]
RD3_NAK_CH=[]
RD3_NAK_SM=[]

RD4_NAK_KH=[]
RD4_NAK_KM=[]
RD4_NAK_YZ=[]
RD4_NAK_CH=[]
RD4_NAK_SM=[]

RD_NAK_KH=[]
RD_NAK_KM=[]
RD_NAK_YZ=[]
RD_NAK_CH=[]
RD_NAK_SM=[]





PS_Traffic_2G_KH=[]
PS_Traffic_2G_KM=[]
PS_Traffic_2G_YZ=[]
PS_Traffic_2G_CH=[]
PS_Traffic_2G_SM=[]


PS_Traffic_3G_KH=[]
PS_Traffic_3G_KM=[]
PS_Traffic_3G_YZ=[]
PS_Traffic_3G_CH=[]
PS_Traffic_3G_SM=[]



PS_Traffic_4G_KH=[]
PS_Traffic_4G_KM=[]
PS_Traffic_4G_YZ=[]
PS_Traffic_4G_CH=[]
PS_Traffic_4G_SM=[]




Week_Vec=[]

for i in range(len(RD_Province_Table)):
    Row_Data=str(RD_Province_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Province=Row_Data[1]
    Contractor=Row_Data[2]
    RD2_Str=Row_Data[3]
    RD2_Val=round(float(RD2_Str[0:len(RD2_Str)-1]),2)
    RD3_Str=Row_Data[4]
    RD3_Val=round(float(RD3_Str[0:len(RD3_Str)-1]),2)
    RD4_Str=Row_Data[5]
    RD4_Val=round(float(RD4_Str[0:len(RD4_Str)-1]),2)
    RD_Str=Row_Data[6]
    RD_Val=round(float(RD_Str[0:len(RD_Str)-1]),2)
    PS_2G_Str=Row_Data[7]
    PS_2G_Val=round(float(PS_2G_Str[0:len(PS_2G_Str)-1]),2)
    PS_3G_Str=Row_Data[8]
    PS_3G_Val=round(float(PS_3G_Str[0:len(PS_3G_Str)-1]),2)
    PS_4G_Str=Row_Data[9]
    PS_4G_Val=round(float(PS_4G_Str[0:len(PS_4G_Str)-1]),2)
    PS_Str=Row_Data[10]
    PS_Val=round(float(PS_Str[0:len(PS_Str)-1]),2)

    Week=Week[2:9]
    Contractor=Contractor[1:len(Contractor)-1]
    Province=Province[1:len(Province)-1]

    #if Week=='1401-33':
    #    break

    #if (Contractor=='NAK-Alborz'):
    #    Week_Vec.append('W'+Week[5:7])
    #    RD2_NAK_Alborz.append(RD2_Val)
    #    RD3_NAK_Alborz.append(RD3_Val)
    #    RD_NAK_Alborz.append(RD_Val)
    #    PS_Traffic_2G_NAK_Alborz.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_NAK_Alborz.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_NAK_Alborz.append(PS_4G_Val/1e3)
    #if (Contractor=='NAK-Tehran'):
    #    RD2_NAK_Tehran.append(RD2_Val)
    #    RD3_NAK_Tehran.append(RD3_Val)
    #    RD_NAK_Tehran.append(RD_Val)
    #    PS_Traffic_2G_NAK_Tehran.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_NAK_Tehran.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_NAK_Tehran.append(PS_4G_Val/1e3)
    #if (Contractor=='NAK-North'):
    #    RD2_NAK_North.append(RD2_Val)
    #    RD3_NAK_North.append(RD3_Val)
    #    RD_NAK_North.append(RD_Val)
    #    PS_Traffic_2G_NAK_North.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_NAK_North.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_NAK_North.append(PS_4G_Val/1e3)
    if (Contractor=='NAK-Nokia' and Province=='KH'):
        Week_Vec.append('W'+Week[5:7])
        RD2_NAK_KH.append(RD2_Val)
        RD3_NAK_KH.append(RD3_Val)
        RD4_NAK_KH.append(RD4_Val)
        RD_NAK_KH.append(RD_Val)
        PS_Traffic_2G_KH.append(PS_2G_Val)
        PS_Traffic_3G_KH.append(PS_3G_Val)
        PS_Traffic_4G_KH.append(PS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='KM'):
        RD2_NAK_KM.append(RD2_Val)
        RD3_NAK_KM.append(RD3_Val)
        RD4_NAK_KM.append(RD4_Val)
        RD_NAK_KM.append(RD_Val)
        PS_Traffic_2G_KM.append(PS_2G_Val)
        PS_Traffic_3G_KM.append(PS_3G_Val)
        PS_Traffic_4G_KM.append(PS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='YZ'):
        RD2_NAK_YZ.append(RD2_Val)
        RD3_NAK_YZ.append(RD3_Val)
        RD4_NAK_YZ.append(RD4_Val)
        RD_NAK_YZ.append(RD_Val)
        PS_Traffic_2G_YZ.append(PS_2G_Val)
        PS_Traffic_3G_YZ.append(PS_3G_Val)
        PS_Traffic_4G_YZ.append(PS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='CH'):
        RD2_NAK_CH.append(RD2_Val)
        RD3_NAK_CH.append(RD3_Val)
        RD4_NAK_CH.append(RD4_Val)
        RD_NAK_CH.append(RD_Val)
        PS_Traffic_2G_CH.append(PS_2G_Val)
        PS_Traffic_3G_CH.append(PS_3G_Val)
        PS_Traffic_4G_CH.append(PS_4G_Val)
    if (Contractor=='NAK-Nokia' and Province=='SM'):
        RD2_NAK_SM.append(RD2_Val)
        RD3_NAK_SM.append(RD3_Val)
        RD4_NAK_SM.append(RD4_Val)
        RD_NAK_SM.append(RD_Val)
        PS_Traffic_2G_SM.append(PS_2G_Val)
        PS_Traffic_3G_SM.append(PS_3G_Val)
        PS_Traffic_4G_SM.append(PS_4G_Val)
    #if (Contractor=='NAK-Huawei'):
    #    RD2_NAK_Huawei.append(RD2_Val)
    #    RD3_NAK_Huawei.append(RD3_Val)
    #    RD_NAK_Huawei.append(RD_Val)
    #    PS_Traffic_2G_NAK_Huawei.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_NAK_Huawei.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_NAK_Huawei.append(PS_4G_Val/1e3)
    #if (Contractor=='Farafan'):
    #    RD2_Farafan.append(RD2_Val)
    #    RD3_Farafan.append(RD3_Val)
    #    RD_Farafan.append(RD_Val)
    #    PS_Traffic_2G_Farafan.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_Farafan.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_Farafan.append(PS_4G_Val/1e3)
    #if (Contractor=='BR-TEL'):
    #    RD2_BR_TEL.append(RD2_Val)
    #    RD3_BR_TEL.append(RD3_Val)
    #    RD_BR_TEL.append(RD_Val)
    #    PS_Traffic_2G_BR_TEL.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_BR_TEL.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_BR_TEL.append(PS_4G_Val/1e3)
    #if (Contractor=='Huawei'):
    #    RD2_Huawei.append(RD2_Val)
    #    RD3_Huawei.append(RD3_Val)
    #    RD_Huawei.append(RD_Val)
    #    PS_Traffic_2G_Huawei.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_Huawei.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_Huawei.append(PS_4G_Val/1e3)
    #if (Contractor=='IRAN'):
    #    RD2_Iran.append(RD2_Val)
    #    RD3_Iran.append(RD3_Val)
    #    RD_Iran.append(RD_Val)
    #    PS_Traffic_2G_Iran.append(PS_2G_Val/1e3)
    #    PS_Traffic_3G_Iran.append(PS_3G_Val/1e3)
    #    PS_Traffic_4G_Iran.append(PS_4G_Val/1e3)






for t in range(8):
        #if (t==0 ):
        #    Contractor="NAK-Alborz"
        #    L2100_Hourly=L2100_Hourly_NAK_Alborz
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Alborz
        #if (t==1 ):
        #    Contractor="NAK-North"
        #    L2100_Hourly=L2100_Hourly_NAK_North
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_North
        if (t==2 ):
            Contractor="NAK-Tehran"

            #downsample_Rate=round(len(L2100Time_Hourly_TH)/50)
            #fig, ax1 = plt.subplots(figsize=(cm_to_inch(35),cm_to_inch(17)))
            #x_Downsample=downsample(L2100Time_Hourly_TH,downsample_Rate)
            #X_Vec=[]
            #x_index=0
            #while len(X_Vec)!=len(x_Downsample):
            #     X_Vec.append(x_index)
            #     x_index=x_index+downsample_Rate
            #ax1.plot(L2100Time_Hourly_KH, L2100_Hourly_TH, label='Tehran')
            #ax1.set_xticks(X_Vec, x_Downsample,fontsize=7, rotation='vertical')
            #leg = ax1.legend();
            #font1 = {'family':'serif','color':'black','size':12}
            #plt.title(Contractor+" L2100 Payload (TB)", fontdict = font1)
            #grid(True)
            #plt.savefig("L2100_"+Contractor+".png")

            #image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\L2100_"+Contractor+".png")
            #y=100
            #x=20
            #h=1300
            #w=1050
            #L2100_Cropped = image[x:w, y:h]
            #cv2.imwrite("L2100_"+Contractor+".png", L2100_Cropped)


        #if (t==3 ):
        #    Contractor="NAK-Huawei"
        #    L2100_Hourly=L2100_Hourly_NAK_Huawei
        #    L2100Time_Hourly=L2100Time_Hourly_NAK_Huawei
        if (t==4 ):
            Contractor="NAK-Nokia"

            Last_RD2_NAK_Nokia=[RD2_NAK_KH[len(RD2_NAK_KH)-1], RD2_NAK_KM[len(RD2_NAK_KM)-1], RD2_NAK_YZ[len(RD2_NAK_YZ)-1], RD2_NAK_CH[len(RD2_NAK_CH)-1], RD2_NAK_SM[len(RD2_NAK_SM)-1]]
            Index_of_Sort_RD2=np.argsort(Last_RD2_NAK_Nokia)
            Data_Sorted_Array_RD2=[]
            x_Labels_RD2=[];
            for k in range(len(Index_of_Sort_RD2)):
                if Index_of_Sort_RD2[k]==0:
                    Data_Sorted_Array_RD2.append(RD2_NAK_KH)
                    x_Labels_RD2.append('Khorasan Razavi')
                if Index_of_Sort_RD2[k]==1:
                    Data_Sorted_Array_RD2.append(RD2_NAK_KM)
                    x_Labels_RD2.append('Kerman')
                if Index_of_Sort_RD2[k]==2:
                    Data_Sorted_Array_RD2.append(RD2_NAK_YZ)
                    x_Labels_RD2.append('Yazd')
                if Index_of_Sort_RD2[k]==3:
                    Data_Sorted_Array_RD2.append(RD2_NAK_CH)
                    x_Labels_RD2.append('Chaharmahal')
                if Index_of_Sort_RD2[k]==4:
                    Data_Sorted_Array_RD2.append(RD2_NAK_SM)
                    x_Labels_RD2.append('Semnan')


            data=np.array(Data_Sorted_Array_RD2)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_RD2_NAK_Nokia):
                plt.text( i + dx[31],Last_RD2_NAK_Nokia[Index_of_Sort_RD2[i]] , str(Last_RD2_NAK_Nokia[Index_of_Sort_RD2[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels_RD2)
            font1 = {'family':'serif','color':'black','size':17}
            plt.title("RD2(%)", fontdict = font1)
            plt.ylim(50, 100)
            grid(True)
            plt.savefig("RD2_"+Contractor+"_Provinces.png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD2_"+Contractor+"_Provinces.png")
            y=80
            x=10
            h=1000
            w=520
            RD2_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite("RD2_"+Contractor+"_Provinces.png", RD2_Bar_Cropped)





            Last_RD3_NAK_Nokia=[RD3_NAK_KH[len(RD3_NAK_KH)-1], RD3_NAK_KM[len(RD3_NAK_KM)-1], RD3_NAK_YZ[len(RD3_NAK_YZ)-1], RD3_NAK_CH[len(RD3_NAK_CH)-1], RD3_NAK_SM[len(RD3_NAK_SM)-1]]
            Index_of_Sort_RD3=np.argsort(Last_RD3_NAK_Nokia)
            Data_Sorted_Array_RD3=[]
            x_Labels_RD3=[];
            for k in range(len(Index_of_Sort_RD3)):
                if Index_of_Sort_RD3[k]==0:
                    Data_Sorted_Array_RD3.append(RD3_NAK_KH)
                    x_Labels_RD3.append('Khorasan Razavi')
                if Index_of_Sort_RD3[k]==1:
                    Data_Sorted_Array_RD3.append(RD3_NAK_KM)
                    x_Labels_RD3.append('Kerman')
                if Index_of_Sort_RD3[k]==2:
                    Data_Sorted_Array_RD3.append(RD3_NAK_YZ)
                    x_Labels_RD3.append('Yazd')
                if Index_of_Sort_RD3[k]==3:
                    Data_Sorted_Array_RD3.append(RD3_NAK_CH)
                    x_Labels_RD3.append('Chaharmahal')
                if Index_of_Sort_RD3[k]==4:
                    Data_Sorted_Array_RD3.append(RD3_NAK_SM)
                    x_Labels_RD3.append('Semnan')


            data=np.array(Data_Sorted_Array_RD3)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_RD3_NAK_Nokia):
                plt.text( i + dx[31],Last_RD3_NAK_Nokia[Index_of_Sort_RD3[i]] , str(Last_RD3_NAK_Nokia[Index_of_Sort_RD3[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels_RD3)
            font1 = {'family':'serif','color':'black','size':17}
            plt.title("RD3(%)", fontdict = font1)
            plt.ylim(75, 100)
            grid(True)
            plt.savefig("RD3_"+Contractor+"_Provinces.png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD3_"+Contractor+"_Provinces.png")
            y=80
            x=10
            h=1000
            w=520
            RD3_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite("RD3_"+Contractor+"_Provinces.png", RD3_Bar_Cropped)




            Last_RD4_NAK_Nokia=[RD4_NAK_KH[len(RD4_NAK_KH)-1], RD4_NAK_KM[len(RD4_NAK_KM)-1], RD4_NAK_YZ[len(RD4_NAK_YZ)-1], RD4_NAK_CH[len(RD4_NAK_CH)-1], RD4_NAK_SM[len(RD4_NAK_SM)-1]]
            Index_of_Sort_RD4=np.argsort(Last_RD4_NAK_Nokia)
            Data_Sorted_Array_RD4=[]
            x_Labels_RD4=[];
            for k in range(len(Index_of_Sort_RD4)):
                if Index_of_Sort_RD4[k]==0:
                    Data_Sorted_Array_RD4.append(RD4_NAK_KH)
                    x_Labels_RD4.append('Khorasan Razavi')
                if Index_of_Sort_RD4[k]==1:
                    Data_Sorted_Array_RD4.append(RD4_NAK_KM)
                    x_Labels_RD4.append('Kerman')
                if Index_of_Sort_RD4[k]==2:
                    Data_Sorted_Array_RD4.append(RD4_NAK_YZ)
                    x_Labels_RD4.append('Yazd')
                if Index_of_Sort_RD4[k]==3:
                    Data_Sorted_Array_RD4.append(RD4_NAK_CH)
                    x_Labels_RD4.append('Chaharmahal')
                if Index_of_Sort_RD4[k]==4:
                    Data_Sorted_Array_RD4.append(RD4_NAK_SM)
                    x_Labels_RD4.append('Semnan')


            data=np.array(Data_Sorted_Array_RD4)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_RD4_NAK_Nokia):
                plt.text( i + dx[31],Last_RD4_NAK_Nokia[Index_of_Sort_RD4[i]] , str(Last_RD4_NAK_Nokia[Index_of_Sort_RD4[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels_RD4)
            font1 = {'family':'serif','color':'black','size':17}
            plt.title("RD4(%)", fontdict = font1)
            plt.ylim(60, 100)
            grid(True)
            plt.savefig("RD4_"+Contractor+"_Provinces.png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD4_"+Contractor+"_Provinces.png")
            y=80
            x=10
            h=1000
            w=520
            RD4_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite("RD4_"+Contractor+"_Provinces.png", RD4_Bar_Cropped)



            Last_RD_NAK_Nokia=[RD_NAK_KH[len(RD_NAK_KH)-1], RD_NAK_KM[len(RD_NAK_KM)-1], RD_NAK_YZ[len(RD_NAK_YZ)-1], RD_NAK_CH[len(RD_NAK_CH)-1], RD_NAK_SM[len(RD_NAK_SM)-1]]
            Index_of_Sort_RD=np.argsort(Last_RD_NAK_Nokia)
            Data_Sorted_Array_RD=[]
            x_Labels_RD=[];
            for k in range(len(Index_of_Sort_RD)):
                if Index_of_Sort_RD[k]==0:
                    Data_Sorted_Array_RD.append(RD_NAK_KH)
                    x_Labels_RD.append('Khorasan Razavi')
                if Index_of_Sort_RD[k]==1:
                    Data_Sorted_Array_RD.append(RD_NAK_KM)
                    x_Labels_RD.append('Kerman')
                if Index_of_Sort_RD[k]==2:
                    Data_Sorted_Array_RD.append(RD_NAK_YZ)
                    x_Labels_RD.append('Yazd')
                if Index_of_Sort_RD[k]==3:
                    Data_Sorted_Array_RD.append(RD_NAK_CH)
                    x_Labels_RD.append('Chaharmahal')
                if Index_of_Sort_RD[k]==4:
                    Data_Sorted_Array_RD.append(RD_NAK_SM)
                    x_Labels_RD.append('Semnan')


            data=np.array(Data_Sorted_Array_RD)

            x = np.arange(data.shape[0])
            dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
            d = 1./(data.shape[1]+2.)

            def cm_to_inch(value):
                return value/2.54
            plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
            axes= plt.axes()


            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


            for i , v in enumerate(Last_RD_NAK_Nokia):
                plt.text( i + dx[31],Last_RD_NAK_Nokia[Index_of_Sort_RD[i]] , str(Last_RD_NAK_Nokia[Index_of_Sort_RD[i]]), color='black', size=12, fontweight='bold')

            axes.set_xticks(x, x_Labels_RD)
            font1 = {'family':'serif','color':'black','size':17}
            plt.title("RD(%)", fontdict = font1)
            plt.ylim(60, 100)
            grid(True)
            plt.savefig("RD_"+Contractor+"_Provinces.png")

            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD_"+Contractor+"_Provinces.png")
            y=80
            x=10
            h=1000
            w=520
            RD_Bar_Cropped = image[x:w, y:h]
            cv2.imwrite("RD_"+Contractor+"_Provinces.png", RD_Bar_Cropped)
 





# ---------------- Total Traffic and Payload per Province and Band


for t in range(25):
        if (t==0 or t==13):
            Contractor="NAK-Alborz"
            Province='KJ'
        if (t==1 or t==14):
            Contractor="NAK-North"
            Province='GL'
        if (t==2 or t==15):
            Contractor="NAK-North"
            Province='GN'
        if (t==3 or t==16):
            Contractor="NAK-North"
            Province='MA'
        if (t==4 or t==17):
            Contractor="NAK-Tehran"
            Province='TH'
        if (t==5 or t==18):
            Contractor="NAK-Huawei"
            Province='AS'
        if (t==6 or t==19):
            Contractor="NAK-Huawei"
            Province='AG'
        if (t==7 or t==20):
            Contractor="NAK-Huawei"
            Province='KZ'
        if (t==8 or t==21):
            Contractor="NAK-Nokia"
            Province='KH'
        if (t==9 or t==22):
            Contractor="NAK-Nokia"
            Province='KM'
        if (t==10 or t==23):
            Contractor="NAK-Nokia"
            Province='YZ'
        if (t==11 or t==24):
            Contractor="NAK-Nokia"
            Province='CH'
        if (t==12 or t==25):
            Contractor="NAK-Nokia"
            Province='SM'            


        Traffic_2G=[]
        Traffic_U900F1=[]
        Traffic_U900F2=[]
        Traffic_U2100F1=[]
        Traffic_U2100F2=[]
        Traffic_U2100F3=[]
        Traffic_L1800F1=[]
        Traffic_L1800F2=[]
        Traffic_L2100F1=[]
        Traffic_L2100F2=[]
        Traffic_L2600F1=[]
        Traffic_L2600F2=[]
        Traffic_L2300F1=[]
        Traffic_L2300F2=[]

        if t<=12:
            conn_performanceDB.execute("select Wk, PIndex, Contractor, sum([2G TCH Traffic]) as '2G TCH Traffic', sum([U900F1 Traffic]) as 'U900F1 Traffic', sum([U900F2 Traffic]) as 'U900F2 Traffic', sum([U2100F1 Traffic]) as 'U2100F1 Traffic', sum([U2100F2 Traffic]) as 'U2100F2 Traffic', sum([U2100F3 Traffic]) as 'U2100F3 Traffic',"+
            " sum([L1800F1 Traffic]) as 'L1800F1 Traffic', sum([L1800F2 Traffic]) as 'L1800F2 Traffic', sum([L2100F1 Traffic]) as 'L2100F1 Traffic', sum([L2100F2 Traffic]) as 'L2100F2 Traffic', sum([L2600F1 Traffic]) as 'L2600F1 Traffic', sum([L2600F2 Traffic]) as 'L2600F2 Traffic', sum([L2300F1 Traffic]) as 'L2300F1 Traffic', sum([L2300F2 Traffic]) as 'L2300F2 Traffic'"+
            " from (select Wk, Contractor,PIndex,  avg([2G TCH Traffic]) as '2G TCH Traffic', avg([U900F1 Traffic]) as 'U900F1 Traffic', avg([U900F2 Traffic]) as 'U900F2 Traffic', avg([U2100F1 Traffic]) as 'U2100F1 Traffic', avg([U2100F2 Traffic]) as 'U2100F2 Traffic', avg([U2100F3 Traffic]) as 'U2100F3 Traffic'," +
            " avg([L1800F1 Traffic]) as 'L1800F1 Traffic', avg([L1800F2 Traffic]) as 'L1800F2 Traffic', avg([L2100F1 Traffic]) as 'L2100F1 Traffic', avg([L2100F2 Traffic]) as 'L2100F2 Traffic', avg([L2600F1 Traffic]) as 'L2600F1 Traffic', avg([L2600F2 Traffic]) as 'L2600F2 Traffic', avg([L2300F1 Traffic]) as 'L2300F1 Traffic', avg([L2300F2 Traffic]) as 'L2300F2 Traffic'"+
            " from Province_KPI_Score_Band_CS_Daily group by Wk, Contractor, PIndex ) tble where Contractor='"+Contractor+"' and PIndex='"+Province+"' group by Wk, PIndex, Contractor  order by Wk")
            Province_Table=conn_performanceDB.fetchall()
        else:
            conn_performanceDB.execute("select Wk,[Province Index], Contractor, sum([2G Payload]) as '2G Payload', sum([U900F1 Payload]) as 'U900F1 Payload', sum([U900F2 Payload]) as 'U900F2 Payload', sum([U2100F1 Payload]) as 'U2100F1 Payload', sum([U2100F2 Payload]) as 'U2100F2 Payload', sum([U2100F3 Payload]) as 'U2100F3 Payload',"+
            " sum([L1800F1 Payload]) as 'L1800F1 Payload', sum([L1800F2 Payload]) as 'L1800F2 Payload', sum([L2100F1 Payload]) as 'L2100F1 Payload', sum([L2100F2 Payload]) as 'L2100F2 Payload', sum([L2600F1 Payload]) as 'L2600F1 Payload', sum([L2600F2 Payload]) as 'L2600F2 Payload', sum([L2300F1 Payload]) as 'L2300F1 Payload', sum([L2300F2 Payload]) as 'L2300F2 Payload'"+
            " from (select Wk, Contractor,[Province Index],  avg([2G PS Traffic (GB)]) as '2G Payload', avg([U900F1 Payload]) as 'U900F1 Payload', avg([U900F2 Payload]) as 'U900F2 Payload', avg([U2100F1 Payload]) as 'U2100F1 Payload', avg([U2100F2 Payload]) as 'U2100F2 Payload', avg([U2100F3 Payload]) as 'U2100F3 Payload'," +
            " avg([L1800F1 Payload]) as 'L1800F1 Payload', avg([L1800F2 Payload]) as 'L1800F2 Payload', avg([L2100F1 Payload]) as 'L2100F1 Payload', avg([L2100F2 Payload]) as 'L2100F2 Payload', avg([L2600F1 Payload]) as 'L2600F1 Payload', avg([L2600F2 Payload]) as 'L2600F2 Payload', avg([L2300F1 Payload]) as 'L2300F1 Payload', avg([L2300F2 Payload]) as 'L2300F2 Payload'"+
            " from Province_KPI_Score_Band_PS_Daily group by Wk, Contractor, [Province Index] ) tble where Contractor='"+Contractor+"' and [Province Index]='"+Province+"' group by Wk, [Province Index], Contractor  order by Wk")
            Province_Table=conn_performanceDB.fetchall()


        for i in range(len(Province_Table)):


            Row_Data=str(Province_Table[i])
            Row_Data=Row_Data.split(", ")

            Week=Row_Data[0]
            Province=Row_Data[1]
            Contractor=Row_Data[2]
            Traffic_2G_Str=Row_Data[3]
            Traffic_2G_Val=round(float(Traffic_2G_Str[0:len(Traffic_2G_Str)-1])/1e3,2)
            Traffic_U900F1_Str=Row_Data[4]
            Traffic_U900F1_Val=round(float(Traffic_U900F1_Str[0:len(Traffic_U900F1_Str)-1])/1e3,2)
            Traffic_U900F2_Str=Row_Data[5]
            Traffic_U900F2_Val=round(float(Traffic_U900F2_Str[0:len(Traffic_U900F2_Str)-1])/1e3,2)
            Traffic_U2100F1_Str=Row_Data[6]
            Traffic_U2100F1_Val=round(float(Traffic_U2100F1_Str[0:len(Traffic_U2100F1_Str)-1])/1e3,2)
            Traffic_U2100F2_Str=Row_Data[7]
            Traffic_U2100F2_Val=round(float(Traffic_U2100F2_Str[0:len(Traffic_U2100F2_Str)-1])/1e3,2)
            Traffic_U2100F3_Str=Row_Data[8]
            Traffic_U2100F3_Val=round(float(Traffic_U2100F3_Str[0:len(Traffic_U2100F3_Str)-1])/1e3,2)
            Traffic_L1800F1_Str=Row_Data[9]
            Traffic_L1800F1_Val=round(float(Traffic_L1800F1_Str[0:len(Traffic_L1800F1_Str)-1])/1e3,2)
            Traffic_L1800F2_Str=Row_Data[10]
            Traffic_L1800F2_Val=round(float(Traffic_L1800F2_Str[0:len(Traffic_L1800F2_Str)-1])/1e3,2)
            Traffic_L2100F1_Str=Row_Data[11]
            Traffic_L2100F1_Val=round(float(Traffic_L2100F1_Str[0:len(Traffic_L2100F1_Str)-1])/1e3,2)
            Traffic_L2100F2_Str=Row_Data[12]
            Traffic_L2100F2_Val=round(float(Traffic_L2100F2_Str[0:len(Traffic_L2100F2_Str)-1])/1e3,2)
            Traffic_L2600F1_Str=Row_Data[13]
            Traffic_L2600F1_Val=round(float(Traffic_L2600F1_Str[0:len(Traffic_L2600F1_Str)-1])/1e3,2)
            Traffic_L2600F2_Str=Row_Data[14]
            Traffic_L2600F2_Val=round(float(Traffic_L2600F2_Str[0:len(Traffic_L2600F2_Str)-1])/1e3,2)
            Traffic_L2300F1_Str=Row_Data[15]
            Traffic_L2300F1_Val=round(float(Traffic_L2300F1_Str[0:len(Traffic_L2300F1_Str)-1])/1e3,2)
            Traffic_L2300F2_Str=Row_Data[16]
            Traffic_L2300F2_Val=round(float(Traffic_L2300F2_Str[0:len(Traffic_L2300F2_Str)-1])/1e3,2)


            Week=Week[2:9]
            Contractor=Contractor[1:len(Contractor)-1]
            Province=Province[1:len(Province)-1]
    

            #if Week=='1401-33':
            #    break

            Traffic_2G.append(Traffic_2G_Val)
            Traffic_U900F1.append(Traffic_U900F1_Val)
            Traffic_U900F2.append(Traffic_U900F2_Val)
            Traffic_U2100F1.append(Traffic_U2100F1_Val)
            Traffic_U2100F2.append(Traffic_U2100F2_Val)
            Traffic_U2100F3.append(Traffic_U2100F3_Val)
            Traffic_L1800F1.append(Traffic_L1800F1_Val)
            Traffic_L1800F2.append(Traffic_L1800F2_Val)
            Traffic_L2100F1.append(Traffic_L2100F1_Val)
            Traffic_L2100F2.append(Traffic_L2100F2_Val)
            Traffic_L2600F1.append(Traffic_L2600F1_Val)
            Traffic_L2600F2.append(Traffic_L2600F2_Val)
            Traffic_L2300F1.append(Traffic_L2300F1_Val)
            Traffic_L2300F2.append(Traffic_L2300F2_Val)


        # Sort Data Based on Last Values
        Last_Traffic_Value=[]
        Last_Traffic_Value=[Traffic_2G[len(Traffic_2G)-1], Traffic_U900F1[len(Traffic_2G)-1], Traffic_U900F2[len(Traffic_2G)-1],  Traffic_U2100F1[len(Traffic_2G)-1], Traffic_U2100F2[len(Traffic_2G)-1], Traffic_U2100F3[len(Traffic_2G)-1], Traffic_L1800F1[len(Traffic_2G)-1],  Traffic_L1800F2[len(Traffic_2G)-1], Traffic_L2100F1[len(Traffic_2G)-1],  Traffic_L2100F2[len(Traffic_2G)-1], Traffic_L2600F1[len(Traffic_2G)-1],  Traffic_L2600F2[len(Traffic_2G)-1], Traffic_L2300F1[len(Traffic_2G)-1],  Traffic_L2300F2[len(Traffic_2G)-1]]
        Index_of_Sort=np.argsort(Last_Traffic_Value)

        Data_Sorted_Array=[]
        x_Labels=[];
        for k in range(len(Index_of_Sort)):
            if Index_of_Sort[k]==0:
                if (np.sum(Traffic_2G)!=0):
                    Data_Sorted_Array.append(Traffic_2G)
                    x_Labels.append('2G')
            if Index_of_Sort[k]==1:
                if (np.sum(Traffic_U900F1)!=0):
                    Data_Sorted_Array.append(Traffic_U900F1)
                    x_Labels.append('U900F1')
            if Index_of_Sort[k]==2:
                if (np.sum(Traffic_U900F2)!=0):
                    Data_Sorted_Array.append(Traffic_U900F2)
                    x_Labels.append('U900F2')
            if Index_of_Sort[k]==3:
                if (np.sum(Traffic_U2100F1)!=0):
                    Data_Sorted_Array.append(Traffic_U2100F1)
                    x_Labels.append('U2100F1')
            if Index_of_Sort[k]==4:
                if (np.sum(Traffic_U2100F2)!=0):
                    Data_Sorted_Array.append(Traffic_U2100F2)
                    x_Labels.append('U2100F2')
            if Index_of_Sort[k]==5:
                if (np.sum(Traffic_U2100F3)!=0):
                    Data_Sorted_Array.append(Traffic_U2100F3)
                    x_Labels.append('U2100F3')
            if Index_of_Sort[k]==6:
                if (np.sum(Traffic_L1800F1)!=0):
                    Data_Sorted_Array.append(Traffic_L1800F1)
                    x_Labels.append('L1800F1')
            if Index_of_Sort[k]==7:
                if (np.sum(Traffic_L1800F2)!=0):
                    Data_Sorted_Array.append(Traffic_L1800F2)
                    x_Labels.append('L1800F2')
            if Index_of_Sort[k]==8:
                if (np.sum(Traffic_L2100F1)!=0):
                    Data_Sorted_Array.append(Traffic_L2100F1)
                    x_Labels.append('L2100F1')
            if Index_of_Sort[k]==9:
                if (np.sum(Traffic_L2100F2)!=0):
                    Data_Sorted_Array.append(Traffic_L2100F2)
                    x_Labels.append('L2100F2')
            if Index_of_Sort[k]==10:
                if (np.sum(Traffic_L2600F1)!=0):
                    Data_Sorted_Array.append(Traffic_L2600F1)
                    x_Labels.append('L2600F1')
            if Index_of_Sort[k]==11:
                if (np.sum(Traffic_L2600F2)!=0):
                    Data_Sorted_Array.append(Traffic_L2600F2)
                    x_Labels.append('L2600F2')
            if Index_of_Sort[k]==12:
                if (np.sum(Traffic_L2300F1)!=0):
                    Data_Sorted_Array.append(Traffic_L2300F1)
                    x_Labels.append('L2300F1')
            if Index_of_Sort[k]==13:
                if (np.sum(Traffic_L2300F2)!=0):
                    Data_Sorted_Array.append(Traffic_L2300F2)
                    x_Labels.append('L2300F2')




        data=np.array(Data_Sorted_Array)

        x = np.arange(data.shape[0])
        dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
        d = 1./(data.shape[1]+2.)

        def cm_to_inch(value):
            return value/2.54
        plt.figure(figsize=(cm_to_inch(28),cm_to_inch(12)))
        axes= plt.axes()

        if t<=12:
            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))

            for k in range(data.shape[0]):
                Last_Value=data[k,data.shape[1]-1]
                plt.text( k+ dx[31],Last_Value , str(Last_Value), color='black', size=10, fontweight='bold')
        else:
            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))

            for k in range(data.shape[0]):
                Last_Value=data[k,data.shape[1]-1]
                plt.text( k+ dx[31],Last_Value , str(Last_Value), color='black', size=10, fontweight='bold')

        #zero_count=0
        #for i , v in enumerate(Last_Traffic_Value):
        #    if Last_Traffic_Value[Index_of_Sort[i]]==0:
        #        zero_count+=1;
        #        continue
        #    plt.text( i -zero_count+ dx[31],Last_Traffic_Value[Index_of_Sort[i]] , str(Last_Traffic_Value[Index_of_Sort[i]]), color='green', size=12, fontweight='bold')

        axes.set_xticks(x, x_Labels,fontsize=7)
        #axes.set_xticks(fontsize=10)
        font1 = {'family':'serif','color':'black','size':14}
        if t<=12:
            plt.title(Province+" Total Traffic (KErlang)", fontdict = font1)
            grid(True)
            plt.savefig('Traffic_'+Province+'_Bar.png')
        else:
            plt.title(Province+" Total Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig('Payload_'+Province+'_Bar.png')



        Data_Sorted_Pie=[]
        for j in range(len(Data_Sorted_Array)):
            Data_Sorted_Pie.append(Data_Sorted_Array[j][len(Data_Sorted_Array[0])-1])

        Data_Sorted_Pie=list([Data_Sorted_Pie]/np.sum(Data_Sorted_Pie)*100)
        fig, ax = plt.subplots(figsize=(8,6))
        x = np.arange(len(Data_Sorted_Pie[0])) 
        width = 0.5
        rects1 = ax.barh(x - width/2, Data_Sorted_Pie[0], width)
        ax.bar_label(rects1,label=Data_Sorted_Pie[0],  fmt='%.2f', fontsize=15)
        ax.set_yticks(x, labels=x_Labels,fontsize=9)
        if t<=12:
            ax.set_xlabel(Province+' Total Traffic (%)', fontsize=14)
            plt.savefig('Traffic_'+Province+'_Bar_Percentage.png')
        else:
            ax.set_xlabel(Province+' Total Payload (%)', fontsize=14)
            plt.savefig('Payload_'+Province+'_Bar_Percentage.png')

        #slide = prs.slides.add_slide(blank_slide_layout)

        ##pic_width_NAK = int(prs.slide_width *0.1)
        ##pic_left_MCI = int(prs.slide_width *0.9)
        ##pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\\NAK.png", 0, 0, pic_width_NAK)
        ##pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\\MCI.png", pic_left_MCI, 0, pic_width_NAK)

        #if t<=7:
        #    image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Traffic_"+Contractor+"_Bar.png")
        #    y=80
        #    x=20
        #    h=1050
        #    w=520
        #    Traffic_Contractor_Bar = image[x:w, y:h]
        #    cv2.imwrite("Traffic_"+Contractor+"_Bar.png", Traffic_Contractor_Bar)


        #    pic_left_1  = int(prs.slide_width *0.045)
        #    pic_top_1   = 0
        #    pic_width_1 = int(prs.slide_width *0.9)


        #    image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Traffic_"+Contractor+"_Bar_Percentage.png")
        #    y=5
        #    x=40
        #    h=770
        #    w=600
        #    Traffic_Contractor_Bar_Percentage = image[x:w, y:h]
        #    cv2.imwrite("Traffic_"+Contractor+"_Bar_Percentage.png", Traffic_Contractor_Bar_Percentage)
        #    pic_left_3  = int(prs.slide_width *0.24)
        #    pic_top_3   = int(prs.slide_width *0.4)
        #    pic_width_3 = int(prs.slide_width *0.47)


        #    pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Traffic_"+Contractor+"_Bar.png", pic_left_1, pic_top_1, pic_width_1)
        #    pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Traffic_"+Contractor+"_Bar_Percentage.png", pic_left_3, pic_top_3, pic_width_3)

        #    prs.save('test.pptx')
        #else:
        #    image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Payload_"+Contractor+"_Bar.png")
        #    y=80
        #    x=20
        #    h=1050
        #    w=520
        #    Payload_Contractor_Bar = image[x:w, y:h]
        #    cv2.imwrite("Payload_"+Contractor+"_Bar.png", Payload_Contractor_Bar)


        #    pic_left_1  = int(prs.slide_width *0.045)
        #    pic_top_1   = 0
        #    pic_width_1 = int(prs.slide_width *0.9)


        #    image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Payload_"+Contractor+"_Bar_Percentage.png")
        #    y=5
        #    x=40
        #    h=770
        #    w=600
        #    Payload_Contractor_Bar_Percentage = image[x:w, y:h]
        #    cv2.imwrite("Payload_"+Contractor+"_Bar_Percentage.png", Payload_Contractor_Bar_Percentage)
        #    pic_left_3  = int(prs.slide_width *0.24)
        #    pic_top_3   = int(prs.slide_width *0.4)
        #    pic_width_3 = int(prs.slide_width *0.47)


        #    pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Payload_"+Contractor+"_Bar.png", pic_left_1, pic_top_1, pic_width_1)
        #    pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\Payload_"+Contractor+"_Bar_Percentage.png", pic_left_3, pic_top_3, pic_width_3)

        #    prs.save('test.pptx')




# --------------------------- CC and RD -----------


for t in range(5):
        #if (t==0 ):
        #    Contractor="NAK-Alborz"
        #    CS_Traffic=CS_Traffic_NAK_Alborz
        #    CC=CC_NAK_Alborz
        #    PS_Traffic=PS_Traffic_NAK_Alborz
        #    RD=RD_NAK_Alborz
        #    CS_Traffic_2G=CS_Traffic_2G_NAK_Alborz
        #    CC2=CC2_NAK_Alborz
        #    CS_Traffic_3G=CS_Traffic_3G_NAK_Alborz
        #    CC3=CC3_NAK_Alborz
        #    PS_Traffic_2G=PS_Traffic_2G_NAK_Alborz
        #    RD2=RD2_NAK_Alborz
        #    PS_Traffic_3G=PS_Traffic_3G_NAK_Alborz
        #    RD3=RD3_NAK_Alborz
        #    PS_Traffic_4G=PS_Traffic_4G_NAK_Alborz
        #    RD4=RD4_NAK_Alborz
        #if (t==1 ):
        #    Contractor="NAK-North"
        #    CS_Traffic=CS_Traffic_NAK_North
        #    CC=CC_NAK_North
        #    PS_Traffic=PS_Traffic_NAK_North
        #    RD=RD_NAK_North
        #    CS_Traffic_2G=CS_Traffic_2G_NAK_North
        #    CC2=CC2_NAK_North
        #    CS_Traffic_3G=CS_Traffic_3G_NAK_North
        #    CC3=CC3_NAK_North
        #    PS_Traffic_2G=PS_Traffic_2G_NAK_North
        #    RD2=RD2_NAK_North
        #    PS_Traffic_3G=PS_Traffic_3G_NAK_North
        #    RD3=RD3_NAK_North
        #    PS_Traffic_4G=PS_Traffic_4G_NAK_North
        #    RD4=RD4_NAK_North
        #if (t==2 ):
        #    Contractor="NAK-Tehran"
        #    CS_Traffic=CS_Traffic_NAK_Tehran
        #    CC=CC_NAK_Tehran
        #    PS_Traffic=PS_Traffic_NAK_Tehran
        #    RD=RD_NAK_Tehran
        #    CS_Traffic_2G=CS_Traffic_2G_NAK_Tehran
        #    CC2=CC2_NAK_Tehran
        #    CS_Traffic_3G=CS_Traffic_3G_NAK_Tehran
        #    CC3=CC3_NAK_Tehran
        #    PS_Traffic_2G=PS_Traffic_2G_NAK_Tehran
        #    RD2=RD2_NAK_Tehran
        #    PS_Traffic_3G=PS_Traffic_3G_NAK_Tehran
        #    RD3=RD3_NAK_Tehran
        #    PS_Traffic_4G=PS_Traffic_4G_NAK_Tehran
        #    RD4=RD4_NAK_Tehran
        #if (t==3 ):
        #    Contractor="NAK-Huawei"
        #    CS_Traffic=CS_Traffic_NAK_Huawei
        #    CC=CC_NAK_Huawei
        #    PS_Traffic=PS_Traffic_NAK_Huawei
        #    RD=RD_NAK_Huawei
        #    CS_Traffic_2G=CS_Traffic_2G_NAK_Huawei
        #    CC2=CC2_NAK_Huawei
        #    CS_Traffic_3G=CS_Traffic_3G_NAK_Huawei
        #    CC3=CC3_NAK_Huawei
        #    PS_Traffic_2G=PS_Traffic_2G_NAK_Huawei
        #    RD2=RD2_NAK_Huawei
        #    PS_Traffic_3G=PS_Traffic_3G_NAK_Huawei
        #    RD3=RD3_NAK_Huawei
        #    PS_Traffic_4G=PS_Traffic_4G_NAK_Huawei
        #    RD4=RD4_NAK_Huawei
        if (t==0 ):
            Contractor="NAK-Nokia"
            Province="KH"
            CS_Traffic=CS_Traffic_KH
            CC=CC_NAK_KH
            PS_Traffic=PS_Traffic_KH
            RD=RD_NAK_KH
            CS_Traffic_2G=CS_Traffic_2G_KH
            CC2=CC2_NAK_KH
            CS_Traffic_3G=CS_Traffic_3G_KH
            CC3=CC3_NAK_KH
            PS_Traffic_2G=PS_Traffic_2G_KH
            RD2=RD2_NAK_KH
            PS_Traffic_3G=PS_Traffic_3G_KH
            RD3=RD3_NAK_KH
            PS_Traffic_4G=PS_Traffic_4G_KH
            RD4=RD4_NAK_KH
        if (t==1 ):
            Contractor="NAK-Nokia"
            Province="KM"
            CS_Traffic=CS_Traffic_KM
            CC=CC_NAK_KM
            PS_Traffic=PS_Traffic_KM
            RD=RD_NAK_KM
            CS_Traffic_2G=CS_Traffic_2G_KM
            CC2=CC2_NAK_KM
            CS_Traffic_3G=CS_Traffic_3G_KM
            CC3=CC3_NAK_KM
            PS_Traffic_2G=PS_Traffic_2G_KM
            RD2=RD2_NAK_KM
            PS_Traffic_3G=PS_Traffic_3G_KM
            RD3=RD3_NAK_KM
            PS_Traffic_4G=PS_Traffic_4G_KM
            RD4=RD4_NAK_KM
        if (t==2 ):
            Contractor="NAK-Nokia"
            Province="YZ"
            CS_Traffic=CS_Traffic_YZ
            CC=CC_NAK_YZ
            PS_Traffic=PS_Traffic_YZ
            RD=RD_NAK_YZ
            CS_Traffic_2G=CS_Traffic_2G_YZ
            CC2=CC2_NAK_YZ
            CS_Traffic_3G=CS_Traffic_3G_YZ
            CC3=CC3_NAK_YZ
            PS_Traffic_2G=PS_Traffic_2G_YZ
            RD2=RD2_NAK_YZ
            PS_Traffic_3G=PS_Traffic_3G_YZ
            RD3=RD3_NAK_YZ
            PS_Traffic_4G=PS_Traffic_4G_YZ
            RD4=RD4_NAK_YZ
        if (t==3 ):
            Contractor="NAK-Nokia"
            Province="CH"
            CS_Traffic=CS_Traffic_CH
            CC=CC_NAK_CH
            PS_Traffic=PS_Traffic_CH
            RD=RD_NAK_CH
            CS_Traffic_2G=CS_Traffic_2G_CH
            CC2=CC2_NAK_CH
            CS_Traffic_3G=CS_Traffic_3G_CH
            CC3=CC3_NAK_CH
            PS_Traffic_2G=PS_Traffic_2G_CH
            RD2=RD2_NAK_CH
            PS_Traffic_3G=PS_Traffic_3G_CH
            RD3=RD3_NAK_CH
            PS_Traffic_4G=PS_Traffic_4G_CH
            RD4=RD4_NAK_CH
        if (t==4 ):
            Contractor="NAK-Nokia"
            Province="SM"
            CS_Traffic=CS_Traffic_SM
            CC=CC_NAK_SM
            PS_Traffic=PS_Traffic_SM
            RD=RD_NAK_SM
            CS_Traffic_2G=CS_Traffic_2G_SM
            CC2=CC2_NAK_SM
            CS_Traffic_3G=CS_Traffic_3G_SM
            CC3=CC3_NAK_SM
            PS_Traffic_2G=PS_Traffic_2G_SM
            RD2=RD2_NAK_SM
            PS_Traffic_3G=PS_Traffic_3G_SM
            RD3=RD3_NAK_SM
            PS_Traffic_4G=PS_Traffic_4G_SM
            RD4=RD4_NAK_SM
        #if (t==5 ):
        #    Contractor="BR-TEL"
        #    CS_Traffic=CS_Traffic_BR_TEL
        #    CC=CC_BR_TEL
        #    PS_Traffic=PS_Traffic_BR_TEL
        #    RD=RD_BR_TEL
        #    CS_Traffic_2G=CS_Traffic_2G_BR_TEL
        #    CC2=CC2_BR_TEL
        #    CS_Traffic_3G=CS_Traffic_3G_BR_TEL
        #    CC3=CC3_BR_TEL
        #    PS_Traffic_2G=PS_Traffic_2G_BR_TEL
        #    RD2=RD2_BR_TEL
        #    PS_Traffic_3G=PS_Traffic_3G_BR_TEL
        #    RD3=RD3_BR_TEL
        #    PS_Traffic_4G=PS_Traffic_4G_BR_TEL
        #    RD4=RD4_BR_TEL
        #if (t==6 ):
        #    Contractor="Farafan"
        #    CS_Traffic=CS_Traffic_Farafan
        #    CC=CC_Farafan
        #    PS_Traffic=PS_Traffic_Farafan
        #    RD=RD_Farafan
        #    CS_Traffic_2G=CS_Traffic_2G_Farafan
        #    CC2=CC2_Farafan
        #    CS_Traffic_3G=CS_Traffic_3G_Farafan
        #    CC3=CC3_Farafan
        #    PS_Traffic_2G=PS_Traffic_2G_Farafan
        #    RD2=RD2_Farafan
        #    PS_Traffic_3G=PS_Traffic_3G_Farafan
        #    RD3=RD3_Farafan
        #    PS_Traffic_4G=PS_Traffic_4G_Farafan
        #    RD4=RD4_Farafan
        #if (t==7 ):
        #    Contractor="Huawei"
        #    CS_Traffic=CS_Traffic_Huawei
        #    CC=CC_Huawei
        #    PS_Traffic=PS_Traffic_Huawei
        #    RD=RD_Huawei
        #    CS_Traffic_2G=CS_Traffic_2G_Huawei
        #    CC2=CC2_Huawei
        #    CS_Traffic_3G=CS_Traffic_3G_Huawei
        #    CC3=CC3_Huawei
        #    PS_Traffic_2G=PS_Traffic_2G_Huawei
        #    RD2=RD2_Huawei
        #    PS_Traffic_3G=PS_Traffic_3G_Huawei
        #    RD3=RD3_Huawei
        #    PS_Traffic_4G=PS_Traffic_4G_Huawei
        #    RD4=RD4_Huawei

        x = np.arange(len(CS_Traffic))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,CS_Traffic,color = "bisque")
        ax2.plot(Week_Vec,CC,color = "darkorange")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("CC-"+Province, fontdict = font1)
        ax1.legend(['CC(%)    Tatal Traffic (KErlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('CC_'+Province+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC_"+Province+".png")
        y=30
        x=5
        h=1000
        w=550
        CC_Cropped = image[x:w, y:h]
        cv2.imwrite("CC_"+Province+".png", CC_Cropped)




        x = np.arange(len(PS_Traffic))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic,color = "limegreen")
        ax2.plot(Week_Vec,RD,color = "darkgreen")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD-"+Province, fontdict = font1)
        ax1.legend(['RD(%)    Tatal Payload (TB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD_'+Province+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD_"+Province+".png")
        y=30
        x=5
        h=1000
        w=550
        RD_Cropped = image[x:w, y:h]
        cv2.imwrite("RD_"+Province+".png", RD_Cropped)



        x = np.arange(len(CS_Traffic_2G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,CS_Traffic_2G,color = "lightsteelblue")
        ax2.plot(Week_Vec,CC2,color = "royalblue")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("CC2-"+Province, fontdict = font1)
        ax1.legend(['CC2(%)    2G Traffic (Erlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('CC2_'+Province+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC2_"+Province+".png")
        y=30
        x=5
        h=1000
        w=550
        CC2_Cropped = image[x:w, y:h]
        cv2.imwrite("CC2_"+Province+".png", CC2_Cropped)




        x = np.arange(len(CS_Traffic_3G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,CS_Traffic_3G,color = "khaki")
        ax2.plot(Week_Vec,CC3,color = "darkkhaki")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("CC3-"+Province, fontdict = font1)
        ax1.legend(['CC3(%)    3G Traffic (Erlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('CC3_'+Province+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC3_"+Province+".png")
        y=30
        x=5
        h=1000
        w=550
        CC3_Cropped = image[x:w, y:h]
        cv2.imwrite("CC3_"+Province+".png", CC3_Cropped)


        x = np.arange(len(PS_Traffic_2G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic_2G,color = "lightgrey")
        ax2.plot(Week_Vec,RD2,color = "dimgrey")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD2-"+Province, fontdict = font1)
        ax1.legend(['RD2(%)    2G Payload (GB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD2_'+Province+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD2_"+Province+".png")
        y=30
        x=5
        h=1000
        w=550
        RD2_Cropped = image[x:w, y:h]
        cv2.imwrite("RD2_"+Province+".png", RD2_Cropped)





        x = np.arange(len(PS_Traffic_3G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic_3G,color = "thistle")
        ax2.plot(Week_Vec,RD3,color = "purple")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD3-"+Province, fontdict = font1)
        ax1.legend(['RD3(%)    3G Payload (GB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD3_'+Province+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD3_"+Province+".png")
        y=30
        x=5
        h=1000
        w=550
        RD3_Cropped = image[x:w, y:h]
        cv2.imwrite("RD3_"+Province+".png", RD3_Cropped)




        x = np.arange(len(PS_Traffic_4G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic_4G,color = "salmon")
        ax2.plot(Week_Vec,RD4,color = "darkred")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD4-"+Province, fontdict = font1)
        ax1.legend(['RD4(%)    4G Payload (GB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD4_'+Province+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD4_"+Province+".png")
        y=30
        x=5
        h=1000
        w=550
        RD4_Cropped = image[x:w, y:h]
        cv2.imwrite("RD4_"+Province+".png", RD4_Cropped)





        #slide = prs.slides.add_slide(blank_slide_layout)

        #pic_left_1  = int(prs.slide_width *0)
        #pic_top_1   = int(prs.slide_width *0.02)
        #pic_width_1 = int(prs.slide_width *0.5)

        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        #pic_left_1  = int(prs.slide_width *0.5)
        #pic_top_1   = int(prs.slide_width *0.02)
        #pic_width_1 = int(prs.slide_width *0.5)

        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)


        #pic_left_1  = int(prs.slide_width *0)
        #pic_top_1   = int(prs.slide_width *0.29)
        #pic_width_1 = int(prs.slide_width *0.5)

        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC2_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        #pic_left_1  = int(prs.slide_width *0.5)
        #pic_top_1   = int(prs.slide_width *0.29)
        #pic_width_1 = int(prs.slide_width *0.5)

        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\CC3_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)



        #pic_left_1  = int(prs.slide_width *0)
        #pic_top_1   = int(prs.slide_width *0.56)
        #pic_width_1 = int(prs.slide_width *0.33)

        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD2_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        #pic_left_1  = int(prs.slide_width *0.33)
        #pic_top_1   = int(prs.slide_width *0.56)
        #pic_width_1 = int(prs.slide_width *0.33)

        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD3_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        #pic_left_1  = int(prs.slide_width *0.67)
        #pic_top_1   = int(prs.slide_width *0.56)
        #pic_width_1 = int(prs.slide_width *0.33)

        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Regional_Dashboards\Regional_Dashboards\Regional_Dashboards\RD4_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)



        #prs.save('test.pptx')


# ---------------------------------------------------------------------
# ----------------------- Technical Dashboard -------------------------
# ---------------------------------------------------------------------
